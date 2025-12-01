from django.shortcuts import render
import time
from django.utils import timezone
from django.http import HttpResponse, StreamingHttpResponse
import mimetypes
from unidecode import unidecode
from Levenshtein import ratio
from urllib.parse import quote
import requests

# Create your views here.
from django.shortcuts import render, redirect
from django.contrib.auth import login, authenticate, logout
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.urls import reverse
from django.views.decorators.csrf import csrf_protect
from django.views.decorators.cache import never_cache
from django.http import JsonResponse
from django.views.decorators.http import require_http_methods
import json

from .forms import RegisterForm, LoginForm
from .models import User, University, Document, UserActivity, DocumentLike, DocumentView, DocumentDownload, SearchHistory, Course

from django.http import JsonResponse
from django.views.decorators.http import require_POST
from django.views.decorators.csrf import csrf_exempt
import json
@never_cache
def home_login_view(request):
    """Trang chủ với form đăng nhập và đăng ký"""
    if request.user.is_authenticated:
        return redirect('dashboard')  # Chuyển hướng đến dashboard nếu đã đăng nhập
    
    login_form = LoginForm()
    register_form = RegisterForm()
    
    context = {
        'login_form': login_form,
        'register_form': register_form,
        'show_register': request.GET.get('tab') == 'register'
    }
    
    return render(request, 'home/login.html', context)


@csrf_protect
@require_http_methods(["POST"])
def register_view(request):
    """Xử lý đăng ký tài khoản"""
    if request.user.is_authenticated:
        return redirect('dashboard')
    
    form = RegisterForm(request.POST)
    
    if form.is_valid():
        try:
            user = form.save()
            # Tự động đăng nhập sau khi đăng ký thành công
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password1')
            user = authenticate(username=username, password=password)
            
            if user:
                login(request, user)
                messages.success(request, f'Chào mừng {user.get_full_name() or user.username}! Tài khoản của bạn đã được tạo thành công.')
                
                # Trả về JSON response cho AJAX
                if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return JsonResponse({
                        'success': True,
                        'redirect_url': reverse('dashboard'),
                        'message': f'Chào mừng {user.get_full_name() or user.username}!'
                    })
                
                return redirect('dashboard')
            else:
                messages.error(request, 'Đăng ký thành công nhưng không thể đăng nhập tự động. Vui lòng đăng nhập thủ công.')
                
        except Exception as e:
            messages.error(request, f'Có lỗi xảy ra khi tạo tài khoản: {str(e)}')
            
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({
                    'success': False,
                    'message': f'Có lỗi xảy ra: {str(e)}'
                })
    else:
        # Xử lý lỗi form
        error_messages = []
        for field, errors in form.errors.items():
            for error in errors:
                if field == '__all__':
                    error_messages.append(error)
                else:
                    field_name = form.fields[field].label or field
                    error_messages.append(f'{field_name}: {error}')
        
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            return JsonResponse({
                'success': False,
                'errors': form.errors,
                'message': '; '.join(error_messages)
            })
        
        for msg in error_messages:
            messages.error(request, msg)
    
    # Nếu có lỗi, hiển thị lại form với thông tin đã nhập
    login_form = LoginForm()
    context = {
        'login_form': login_form,
        'register_form': form,
        'show_register': True
    }
    
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        return JsonResponse({
            'success': False,
            'errors': form.errors
        })
    
    return render(request, 'home/login.html', context)


@csrf_protect
@require_http_methods(["POST"])
def login_view(request):
    """Xử lý đăng nhập"""
    if request.user.is_authenticated:
        return redirect('dashboard')
    
    form = LoginForm(request.POST)
    
    if form.is_valid():
        username = form.cleaned_data.get('username')
        password = form.cleaned_data.get('password')
        remember_me = form.cleaned_data.get('remember_me', False)
        
        # Xử lý đăng nhập bằng email
        if '@' in username:
            try:
                user_obj = User.objects.get(email=username)
                username = user_obj.username
            except User.DoesNotExist:
                pass
        
        user = authenticate(request, username=username, password=password)
        
        if user is not None:
            if user.is_banned:
                error_msg = f'Tài khoản đã bị cấm. Lý do: {user.ban_reason or "Không rõ lý do"}'
                messages.error(request, error_msg)
                
                if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return JsonResponse({
                        'success': False,
                        'message': error_msg
                    })
            else:
                login(request, user)
                
                # Xử lý remember me
                if not remember_me:
                    request.session.set_expiry(0)  # Hết hạn khi đóng browser
                else:
                    request.session.set_expiry(1209600)  # 2 weeks
                
                messages.success(request, f'Chào mừng trở lại, {user.get_full_name() or user.username}!')
                
                # Chuyển hướng đến trang được yêu cầu hoặc dashboard
                next_url = request.GET.get('next', 'dashboard')
                
                if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return JsonResponse({
                        'success': True,
                        'redirect_url': reverse(next_url) if next_url == 'dashboard' else next_url,
                        'message': f'Chào mừng trở lại, {user.get_full_name() or user.username}!'
                    })
                
                return redirect(next_url)
        else:
            error_msg = 'Tên đăng nhập/Email hoặc mật khẩu không đúng.'
            messages.error(request, error_msg)
            
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({
                    'success': False,
                    'message': error_msg
                })
    else:
        # Xử lý lỗi form
        error_messages = []
        for field, errors in form.errors.items():
            for error in errors:
                if field == '__all__':
                    error_messages.append(error)
                else:
                    field_name = form.fields[field].label or field
                    error_messages.append(f'{field_name}: {error}')
        
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            return JsonResponse({
                'success': False,
                'errors': form.errors,
                'message': '; '.join(error_messages)
            })
        
        for msg in error_messages:
            messages.error(request, msg)
    
    # Nếu đăng nhập thất bại
    register_form = RegisterForm()
    context = {
        'login_form': form,
        'register_form': register_form,
        'show_register': False
    }
    
    return render(request, 'home/login.html', context)


def logout_view(request):
    """Xử lý đăng xuất"""
    if request.user.is_authenticated:
        username = request.user.get_full_name() or request.user.username
        logout(request)
        messages.success(request, f'Tạm biệt {username}! Bạn đã đăng xuất thành công.')
    
    return redirect('dashboard')



from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse
from django.db.models import Sum, Count
from .models import Document, University, Course, UserActivity, ChatRoom

def dashboard_view(request):
    """Trang dashboard - cho phép truy cập không cần đăng nhập"""
    documents = Document.objects.filter(
        status='approved',
        is_public=True
    ).select_related(
        'university', 'course', 'uploaded_by'
    ).order_by('-created_at')[:12]
    
    # Thống kê user - chỉ hiển thị nếu đã đăng nhập
    if request.user.is_authenticated:
        user_documents_count = Document.objects.filter(uploaded_by=request.user).count()
        user_total_downloads = Document.objects.filter(uploaded_by=request.user).aggregate(
            total=Sum('download_count'))['total'] or 0
        user_total_likes = Document.objects.filter(uploaded_by=request.user).aggregate(
            total=Sum('like_count'))['total'] or 0
        
        # Hoạt động gần đây
        recent_activities = UserActivity.objects.filter(
            user=request.user
        ).order_by('-created_at')[:5]
    else:
        user_documents_count = 0
        user_total_downloads = 0
        user_total_likes = 0
        recent_activities = []
    
    # Danh sách trường với thống kê
    universities = University.objects.filter(is_active=True).annotate(
        courses_count=Count('course'),
        documents_count=Count('document')
    ).order_by('name')
    
    student_count = User.objects.filter(is_superuser=False).count()
    
    context = {
        'documents': documents,
        'universities': universities,
        'student_count': student_count,
        'user_documents_count': user_documents_count,
        'user_total_downloads': user_total_downloads,
        'user_total_likes': user_total_likes,
        'recent_activities': recent_activities,
        'public_rooms_count': ChatRoom.objects.filter(room_type='public', is_active=True).count(),
        'online_users_count': 0,  # Logic để đếm user online
        'recent_chat_rooms': ChatRoom.objects.filter(is_active=True).order_by('-updated_at')[:3]
    }
    return render(request, 'dashboard/index.html', context)
from django.db.models import Q
def university_courses_view(request, university_id):
    """API endpoint để lấy danh sách môn học của một trường"""
    university = get_object_or_404(University, id=university_id, is_active=True)
    
    courses = Course.objects.filter(
        university=university,
        is_active=True
    ).annotate(
        documents_count=Count('document', filter=Q(document__status='approved'))
    ).order_by('code')
    
    courses_data = []
    for course in courses:
        courses_data.append({
            'id': course.id,
            'code': course.code,
            'name': course.name,
            'description': course.description,
            'documents_count': course.documents_count
        })
    
    return JsonResponse({
        'university': {
            'id': university.id,
            'name': university.name,
            'short_name': university.short_name
        },
        'courses': courses_data
    })

def course_documents_view(request, course_id):
    """API endpoint để lấy danh sách tài liệu của một môn học"""
    course = get_object_or_404(Course, id=course_id, is_active=True)
    
    documents = Document.objects.filter(
        course=course,
        status='approved',
        is_public=True
    ).select_related(
        'university', 'uploaded_by'
    ).order_by('-created_at')
    
    documents_data = []
    for doc in documents:
        documents_data.append({
            'id': doc.id,
            'title': doc.title,
            'document_type': doc.document_type,
            'document_type_display': doc.get_document_type_display(),
            'thumbnail': doc.thumbnail.url if doc.thumbnail else None,
            'view_count': doc.view_count,
            'download_count': doc.download_count,
            'like_count': doc.like_count,
            'status': doc.status,
            'uploaded_by': doc.uploaded_by.get_full_name() or doc.uploaded_by.username,
            'university_name': doc.university.short_name or doc.university.name,
            'created_at': doc.created_at.strftime('%d/%m/%Y')
        })
    
    return JsonResponse({
        'course': {
            'id': course.id,
            'code': course.code,
            'name': course.name,
            'university': course.university.name
        },
        'documents': documents_data
    })

# Thêm vào urlpatterns trong urls.py:
# path('api/university/<int:university_id>/courses/', views.university_courses_view, name='university_courses'),
# path('api/course/<int:course_id>/documents/', views.course_documents_view, name='course_documents'),
from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from django.http import HttpResponse, Http404
from django.contrib import messages


def documents_search(request):
    query = request.GET.get('q', '').strip()
    documents = Document.objects.filter(status='approved', is_public=True)
    
    if query:
        # Tìm kiếm trong title, description và AI keywords
        from django.db.models import Q
        documents = documents.filter(
            Q(title__icontains=query) |
            Q(description__icontains=query) |
            Q(ai_keywords__overlap=[query])
        ).distinct()
        
        # Save search history
        SearchHistory.objects.create(
            user=request.user if request.user.is_authenticated else None,
            query=query,
            result_count=documents.count(),
            ip_address=request.META.get('REMOTE_ADDR')
        )
    
    # Pagination
    from django.core.paginator import Paginator
    paginator = Paginator(documents, 12)
    page = request.GET.get('page', 1)
    documents = paginator.get_page(page)
    
    context = {
        'documents': documents,
        'search_query': query,
        'total_results': paginator.count,
    }
    
    return render(request, 'documents/list.html', context)

@login_required
def document_like(request, document_id):
    document = get_object_or_404(Document, id=document_id, status='approved', is_public=True)
    
    like, created = DocumentLike.objects.get_or_create(
        document=document,
        user=request.user
    )
    
    if not created:
        # Nếu đã like rồi thì unlike
        like.delete()
        document.like_count = max(0, document.like_count - 1)
        liked = False
    else:
        # Like mới
        document.like_count += 1
        liked = True
    
    document.save(update_fields=['like_count'])
    
    # Ghi log hoạt động
    action = 'like_document' if liked else 'unlike_document'
    UserActivity.objects.create(
        user=request.user,
        action=action,
        description=f'{"Thích" if liked else "Bỏ thích"} tài liệu "{document.title}"',
        document=document,
        ip_address=request.META.get('REMOTE_ADDR'),
        user_agent=request.META.get('HTTP_USER_AGENT', '')
    )
    
    return JsonResponse({
        'liked': liked,
        'like_count': document.like_count
    })

def document_view(request, document_id):
    document = get_object_or_404(Document, id=document_id, status='approved', is_public=True)
    
    # Tăng lượt xem
    document.view_count += 1
    document.save(update_fields=['view_count'])
    
    # Lưu lại lượt xem
    DocumentView.objects.create(
        document=document,
        user=request.user if request.user.is_authenticated else None,
        ip_address=request.META.get('REMOTE_ADDR'),
        user_agent=request.META.get('HTTP_USER_AGENT', '')
    )
    
    # Ghi log hoạt động nếu user đã đăng nhập
    if request.user.is_authenticated:
        UserActivity.objects.create(
            user=request.user,
            action='view_document',
            description=f'Xem tài liệu "{document.title}"',
            document=document,
            ip_address=request.META.get('REMOTE_ADDR'),
            user_agent=request.META.get('HTTP_USER_AGENT', '')
        )
    
    # Kiểm tra user đã like chưa (chỉ khi đã đăng nhập)
    is_liked = False
    if request.user.is_authenticated:
        is_liked = DocumentLike.objects.filter(document=document, user=request.user).exists()
    
    # Tài liệu liên quan
    related_documents = Document.objects.filter(
        course=document.course,
        status='approved',
        is_public=True
    ).exclude(id=document.id)[:6]
    
    context = {
        'document': document,
        'is_liked': is_liked,
        'related_documents': related_documents,
    }
    
    return render(request, 'documents/view.html', context)
from django.utils import timezone

@login_required
def document_download(request, document_id):
    document = get_object_or_404(Document, id=document_id, status='approved', is_public=True)
    print(f"Document ID: {document.id}")
    print(f"File path: {document.file_path}")
    print(f"Has public_id: {hasattr(document.file_path, 'public_id')}")
    if hasattr(document.file_path, 'public_id'):
        print(f"Public ID: {document.file_path.public_id}")
    # Check download restriction for free users
    from .premium_views import check_download_limit, log_download
    
    if not request.user.is_premium:
        can_download, downloads_today = check_download_limit(request.user)
        if not can_download:
            messages.warning(
                request,
                f'Bạn đã tải {downloads_today} tài liệu hôm nay (giới hạn cho tài khoản thường). '
                f'Nâng cấp Premium để tải không giới hạn!'
            )
            return redirect('premium_upgrade')
    
    # Log download
    log_download(request.user, document, request.META.get('REMOTE_ADDR'))
    
    # Tăng lượt tải
    document.download_count += 1
    document.save(update_fields=['download_count'])

    # Lưu lại lượt tải
    DocumentDownload.objects.create(
        document=document,
        user=request.user,
        ip_address=request.META.get('REMOTE_ADDR')
    )
    
    # Ghi log hoạt động 
    UserActivity.objects.create(
        user=request.user,
        action='download_document',
        description=f'Tải tài liệu "{document.title}"',
        document=document,
        ip_address=request.META.get('REMOTE_ADDR'),
        user_agent=request.META.get('HTTP_USER_AGENT', '')
    )   
    
    # Redirect đến file trên Cloudinary
    if document.file_path:
        # Test các cách khác nhau
        original_url = document.file_path.url
        print(f"Original URL: {original_url}")
        
        # Thử redirect trực tiếp trước
        return redirect(original_url)

def check_username_availability(request):
    """API kiểm tra tên đăng nhập có khả dụng không"""
    if request.method == 'GET':
        username = request.GET.get('username', '').strip()
        if len(username) < 3:
            return JsonResponse({
                'available': False,
                'message': 'Tên đăng nhập phải có ít nhất 3 ký tự'
            })
        
        exists = User.objects.filter(username=username).exists()
        return JsonResponse({
            'available': not exists,
            'message': 'Tên đăng nhập đã tồn tại' if exists else 'Tên đăng nhập khả dụng'
        })
    
    return JsonResponse({'available': False, 'message': 'Method not allowed'})


def check_email_availability(request):
    """API kiểm tra email có khả dụng không"""
    if request.method == 'GET':
        email = request.GET.get('email', '').strip()
        if not email:
            return JsonResponse({
                'available': False,
                'message': 'Email không được để trống'
            })
        
        exists = User.objects.filter(email=email).exists()
        return JsonResponse({
            'available': not exists,
            'message': 'Email đã được sử dụng' if exists else 'Email khả dụng'
        })
    
    return JsonResponse({'available': False, 'message': 'Method not allowed'})


# views.py
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import authenticate, login, logout, update_session_auth_hash
from django.contrib.auth.decorators import login_required
from django.contrib.auth.tokens import default_token_generator
from django.contrib.sites.shortcuts import get_current_site
from django.template.loader import render_to_string
from django.utils.http import urlsafe_base64_encode, urlsafe_base64_decode
from django.utils.encoding import force_bytes, force_str
from django.core.mail import EmailMessage
from django.contrib import messages
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_POST
from django.urls import reverse
import json
from .models import User, Document, StudyList, AIQuizAttempt, PremiumTransaction
from .forms import ProfileUpdateForm, PasswordChangeForm, PasswordResetForm, SetPasswordForm
from django.db.models import Sum  # thêm dòng này


@login_required
def profile_view(request):
    """Hiển thị trang profile của user"""
    user = request.user
    
    # Thống kê hoạt động của user
    total_documents = Document.objects.filter(uploaded_by=user).count()
    total_downloads = Document.objects.filter(uploaded_by=user).aggregate(
        total=Sum('download_count')

    )['total'] or 0
    total_likes = Document.objects.filter(uploaded_by=user).aggregate(
        total=Sum('like_count')

    )['total'] or 0
    total_quiz_attempts = AIQuizAttempt.objects.filter(user=user).count()
    
    # Tài liệu gần đây
    recent_documents = Document.objects.filter(
        uploaded_by=user, 
        status='approved'
    ).order_by('-created_at')[:5]
    
    # Study lists
    study_lists = StudyList.objects.filter(user=user).order_by('-updated_at')[:5]
    
    # Premium info
    premium_transaction = PremiumTransaction.objects.filter(
        user=user, 
        status='completed'
    ).order_by('-created_at').first()
    
    context = {
        'user': user,
        'total_documents': total_documents,
        'total_downloads': total_downloads,
        'total_likes': total_likes,
        'total_quiz_attempts': total_quiz_attempts,
        'recent_documents': recent_documents,
        'study_lists': study_lists,
        'premium_transaction': premium_transaction,
    }
    
    return render(request, 'accounts/profile.html', context)


@login_required
def profile_edit(request):
    """Chỉnh sửa thông tin profile"""
    if request.method == 'POST':
        form = ProfileUpdateForm(request.POST, instance=request.user)
        if form.is_valid():
            form.save()
            messages.success(request, 'Thông tin profile đã được cập nhật thành công!')
            return redirect('profile')
    else:
        form = ProfileUpdateForm(instance=request.user)
    
    return render(request, 'accounts/profile_edit.html', {'form': form})


@login_required
def change_password(request):
    """Đổi mật khẩu cho user đã đăng nhập"""
    if request.method == 'POST':
        form = PasswordChangeForm(user=request.user, data=request.POST)
        if form.is_valid():
            user = form.save()
            update_session_auth_hash(request, user)  # Giữ session sau khi đổi pass
            messages.success(request, 'Mật khẩu đã được thay đổi thành công!')
            return redirect('profile')
        else:
            messages.error(request, 'Vui lòng kiểm tra lại thông tin.')
    else:
        form = PasswordChangeForm(user=request.user)
    
    return render(request, 'accounts/change_password.html', {'form': form})


def forgot_password(request):
    """Gửi email reset mật khẩu"""
    if request.method == 'POST':
        form = PasswordResetForm(request.POST)
        if form.is_valid():
            email = form.cleaned_data['email']
            
            # Tìm user theo email
            try:
                user = User.objects.get(email=email)
            except User.DoesNotExist:
                messages.error(request, 'Email này không tồn tại trong hệ thống.')
                return render(request, 'accounts/forgot_password.html', {'form': form})
            
            # Tạo token reset
            token = default_token_generator.make_token(user)
            uid = urlsafe_base64_encode(force_bytes(user.pk))
            
            # Tạo link reset
            current_site = get_current_site(request)
            reset_link = f"http://{current_site.domain}/reset-password/{uid}/{token}/"
            
            # Gửi email
            subject = 'StudyBot - Đặt lại mật khẩu'
            message = render_to_string('accounts/password_reset_email.html', {
                'user': user,
                'reset_link': reset_link,
                'site_name': 'StudyBot'
            })
            
            email_msg = EmailMessage(subject, message, to=[email])
            email_msg.content_subtype = 'html'
            
            try:
                email_msg.send()
                messages.success(request, 'Email đặt lại mật khẩu đã được gửi! Vui lòng kiểm tra hộp thư.')
                return redirect('home_login')
            except Exception as e:
                messages.error(request, 'Có lỗi xảy ra khi gửi email. Vui lòng thử lại.')
    else:
        form = PasswordResetForm()
    
    return render(request, 'accounts/forgot_password.html', {'form': form})


def reset_password(request, uidb64, token):
    """Reset mật khẩu từ link trong email"""
    try:
        uid = force_str(urlsafe_base64_decode(uidb64))
        user = User.objects.get(pk=uid)
    except (TypeError, ValueError, OverflowError, User.DoesNotExist):
        user = None
    
    if user is not None and default_token_generator.check_token(user, token):
        if request.method == 'POST':
            form = SetPasswordForm(user=user, data=request.POST)
            if form.is_valid():
                form.save()
                messages.success(request, 'Mật khẩu đã được đặt lại thành công! Vui lòng đăng nhập.')
                return redirect('login')
        else:
            form = SetPasswordForm(user=user)
        
        return render(request, 'accounts/reset_password.html', {
            'form': form,
            'valid_link': True
        })
    else:
        return render(request, 'accounts/reset_password.html', {
            'valid_link': False
        })
import uuid
import cloudinary.uploader

@login_required
@require_POST
def upload_avatar(request):
    """Upload avatar lên Cloudinary (AJAX)"""
    if 'avatar' in request.FILES:
        avatar_file = request.FILES['avatar']
        
        # Validate file size
        if avatar_file.size > 5 * 1024 * 1024:  # 5MB
            return JsonResponse({'success': False, 'error': 'File quá lớn (tối đa 5MB)'})
        
        # Validate file type
        if not avatar_file.content_type.startswith('image/'):
            return JsonResponse({'success': False, 'error': 'Chỉ được upload file ảnh'})
        
        try:
            # Import cloudinary uploader
            import cloudinary.uploader
            
            # Upload file lên Cloudinary
            # Cloudinary sẽ tự động tạo public_id unique nếu không chỉ định
            upload_result = cloudinary.uploader.upload(
                avatar_file,
                folder="avatars/",
                public_id=f"user_{request.user.id}_{uuid.uuid4().hex[:8]}",  # Tạo ID unique
                overwrite=True,
                resource_type="image",
                transformation=[
                    {'width': 300, 'height': 300, 'crop': 'fill', 'gravity': 'face'},  # Resize và crop
                    {'quality': 'auto:good'}  # Tối ưu chất lượng
                ]
            )
            
            # Lưu CloudinaryField sẽ tự động lưu public_id
            request.user.avatar = upload_result['public_id']
            request.user.save(update_fields=['avatar'])
            
            # Ghi log hoạt động
            UserActivity.objects.create(
                user=request.user,
                action='update_avatar',
                description='Cập nhật ảnh đại diện',
                ip_address=request.META.get('REMOTE_ADDR'),
                user_agent=request.META.get('HTTP_USER_AGENT', '')
            )
            
            return JsonResponse({
                'success': True, 
                'avatar_url': request.user.avatar.url,  # CloudinaryField tự động tạo URL
                'message': 'Cập nhật ảnh đại diện thành công!'
            })
            
        except Exception as e:
            # Log lỗi để debug
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"Upload avatar failed for user {request.user.id}: {str(e)}")
            
            return JsonResponse({
                'success': False, 
                'error': f'Lỗi upload ảnh: {str(e)}'
            })
    
    return JsonResponse({'success': False, 'error': 'Không có file được chọn'})


@login_required
@require_POST  
def delete_avatar(request):
    """Xóa avatar (AJAX)"""
    try:
        if request.user.avatar:
            # Xóa file trên Cloudinary
            import cloudinary.uploader
            cloudinary.uploader.destroy(request.user.avatar.public_id)
            
            # Xóa trong database
            request.user.avatar = None
            request.user.save(update_fields=['avatar'])
            
            # Ghi log
            UserActivity.objects.create(
                user=request.user,
                action='delete_avatar',
                description='Xóa ảnh đại diện',
                ip_address=request.META.get('REMOTE_ADDR'),
                user_agent=request.META.get('HTTP_USER_AGENT', '')
            )
            
            return JsonResponse({
                'success': True,
                'message': 'Đã xóa ảnh đại diện'
            })
        else:
            return JsonResponse({
                'success': False,
                'error': 'Không có ảnh đại diện để xóa'
            })
            
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': f'Lỗi khi xóa ảnh: {str(e)}'
        })
    

# chat/views.py - Thêm vào file views.py hiện tại

from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.http import JsonResponse
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import csrf_exempt
from django.core.paginator import Paginator
from django.db.models import Q, Count, Max
from django.utils import timezone
from django.db import transaction
from .models import *
import json


def chat_rooms_list(request):
    """Danh sách phòng chat"""
    # Filter parameters
    room_type = request.GET.get('room_type', '')
    university_id = request.GET.get('university', '')
    course_id = request.GET.get('course', '')
    search = request.GET.get('search', '')
    
    # Base queryset
    rooms = ChatRoom.objects.filter(is_active=True).select_related(
        'created_by', 'university', 'course'
    ).annotate(
        members_count=Count('chatroommember', distinct=True),  # Thêm distinct=True
        last_message_time=Max('chatmessage__created_at')
    ).order_by('-last_message_time', '-created_at')
    
    # Apply filters
    if room_type:
        rooms = rooms.filter(room_type=room_type)
    
    if university_id:
        rooms = rooms.filter(university_id=university_id)
        
    if course_id:
        rooms = rooms.filter(course_id=course_id)
        
    if search:
        rooms = rooms.filter(
            Q(name__icontains=search) |
            Q(description__icontains=search)
        )
    
    # Pagination
    paginator = Paginator(rooms, 12)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)
    
    # Get filter data
    universities = University.objects.filter(is_active=True)
    courses = Course.objects.filter(is_active=True).select_related('university')
    
    # User's joined rooms (chỉ khi đã đăng nhập)
    user_rooms = []
    if request.user.is_authenticated:
        user_rooms = ChatRoomMember.objects.filter(user=request.user).values_list('room_id', flat=True)
    
    context = {
        'page_obj': page_obj,
        'universities': universities,
        'courses': courses,
        'user_rooms': user_rooms,
        'current_filters': {
            'room_type': room_type,
            'university': university_id,
            'course': course_id,
            'search': search,
        }
    }
    
    return render(request, 'chat/rooms_list.html', context)


def chat_room_detail(request, room_id):
    """Chi tiết phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check if user is authenticated and member
    is_member = False
    if request.user.is_authenticated:
        is_member = ChatRoomMember.objects.filter(room=room, user=request.user).exists()
    
    # If user is not authenticated, redirect to login
    if not request.user.is_authenticated:
        messages.info(request, f'Vui lòng đăng nhập để tham gia phòng "{room.name}"')
        return redirect('home_login')
    
    # If room is private and user is not member, check password
    if room.room_type == 'private' and not is_member:
        if request.method == 'POST':
            password = request.POST.get('password')
            if room.password and room.password == password:
                # Join room
                ChatRoomMember.objects.get_or_create(
                    room=room,
                    user=request.user,
                    defaults={'role': 'member'}
                )
                is_member = True
                messages.success(request, f'Bạn đã tham gia phòng "{room.name}"!')
            else:
                messages.error(request, 'Mật khẩu không chính xác!')
                return render(request, 'chat/room_password.html', {'room': room})
        else:
            return render(request, 'chat/room_password.html', {'room': room})
    
    # If not member of public room, auto join
    if not is_member and room.room_type == 'public':
        # Check room capacity
        members_count = ChatRoomMember.objects.filter(room=room).count()
        if members_count < room.max_members:
            ChatRoomMember.objects.get_or_create(
                room=room,
                user=request.user,
                defaults={'role': 'member'}
            )
            is_member = True
        else:
            messages.error(request, 'Phòng đã đầy! Không thể tham gia.')
            return redirect('chat_rooms_list')
    
    if not is_member:
        messages.error(request, 'Bạn không có quyền truy cập phòng này!')
        return redirect('chat_rooms_list')
    
    # Get user's membership info
    membership = ChatRoomMember.objects.get(room=room, user=request.user)
    
    # Update last seen
    membership.last_seen = timezone.now()
    membership.save()
    
    # Get messages (latest 50)
    messages_list = ChatMessage.objects.filter(
        room=room, 
        is_deleted=False
    ).select_related(
        'user', 'reply_to__user'
    ).order_by('-created_at')[:50]
    
    messages_list = list(reversed(messages_list))
    
    # Get room members
    members = ChatRoomMember.objects.filter(room=room).select_related('user').order_by(
        '-role', '-last_seen'
    )
    
    context = {
        'room': room,
        'membership': membership,
        'messages': messages_list,
        'members': members,
        'is_admin': membership.role in ['admin', 'moderator'],
    }
    
    return render(request, 'chat/room_detail.html', context)


@login_required
def chat_room_create(request):
    """Tạo phòng chat mới"""
    if request.method == 'POST':
        name = request.POST.get('name', '').strip()
        description = request.POST.get('description', '').strip()
        room_type = request.POST.get('room_type', 'public')
        password = request.POST.get('password', '').strip()
        max_members = int(request.POST.get('max_members', 100))
        university_id = request.POST.get('university')
        course_id = request.POST.get('course')
        
        # Validation
        if not name:
            messages.error(request, 'Vui lòng nhập tên phòng!')
            return render(request, 'chat/room_create.html', {
                'universities': University.objects.filter(is_active=True),
                'courses': Course.objects.filter(is_active=True).select_related('university')
            })
        
        if room_type == 'private' and not password:
            messages.error(request, 'Phòng riêng tư cần có mật khẩu!')
            return render(request, 'chat/room_create.html', {
                'universities': University.objects.filter(is_active=True),
                'courses': Course.objects.filter(is_active=True).select_related('university')
            })
        
        with transaction.atomic():
            # Create room
            room = ChatRoom.objects.create(
                name=name,
                description=description,
                room_type=room_type,
                password=password if room_type == 'private' else None,
                max_members=max_members,
                created_by=request.user,
                university_id=university_id if university_id else None,
                course_id=course_id if course_id else None,
            )
            
            # Add creator as admin
            ChatRoomMember.objects.create(
                room=room,
                user=request.user,
                role='admin'
            )
            
            # Log activity
            UserActivity.objects.create(
                user=request.user,
                action='create_chat_room',
                description=f'Tạo phòng chat "{name}"',
                chat_room=room
            )
        
        messages.success(request, f'Tạo phòng "{name}" thành công!')
        return redirect('chat_room_detail', room_id=room.id)
    
    universities = University.objects.filter(is_active=True)
    courses = Course.objects.filter(is_active=True).select_related('university')
    
    context = {
        'universities': universities,
        'courses': courses,
    }
    
    return render(request, 'chat/room_create.html', context)


@login_required
def chat_room_edit(request, room_id):
    """Chỉnh sửa phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check permissions
    membership = ChatRoomMember.objects.filter(room=room, user=request.user).first()
    if not membership or membership.role not in ['admin', 'moderator']:
        messages.error(request, 'Bạn không có quyền chỉnh sửa phòng này!')
        return redirect('chat_room_detail', room_id=room.id)
    
    if request.method == 'POST':
        name = request.POST.get('name', '').strip()
        description = request.POST.get('description', '').strip()
        room_type = request.POST.get('room_type', room.room_type)
        password = request.POST.get('password', '').strip()
        max_members = int(request.POST.get('max_members', room.max_members))
        university_id = request.POST.get('university')
        course_id = request.POST.get('course')
        
        if not name:
            messages.error(request, 'Vui lòng nhập tên phòng!')
        elif room_type == 'private' and not password and not room.password:
            messages.error(request, 'Phòng riêng tư cần có mật khẩu!')
        else:
            # Update room
            room.name = name
            room.description = description
            room.room_type = room_type
            room.max_members = max_members
            room.university_id = university_id if university_id else None
            room.course_id = course_id if course_id else None
            
            # Update password only if provided
            if password:
                room.password = password
            elif room_type != 'private':
                room.password = None
                
            room.save()
            
            messages.success(request, 'Cập nhật phòng thành công!')
            return redirect('chat_room_detail', room_id=room.id)
    
    universities = University.objects.filter(is_active=True)
    courses = Course.objects.filter(is_active=True).select_related('university')
    
    context = {
        'room': room,
        'universities': universities,
        'courses': courses,
    }
    
    return render(request, 'chat/room_edit.html', context)


@login_required
@require_http_methods(["POST"])
def chat_room_delete(request, room_id):
    """Xóa phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Only creator can delete
    if room.created_by != request.user:
        messages.error(request, 'Bạn không có quyền xóa phòng này!')
        return redirect('chat_room_detail', room_id=room.id)
    
    room_name = room.name
    room.is_active = False
    room.save()
    
    messages.success(request, f'Đã xóa phòng "{room_name}"!')
    return redirect('chat_rooms_list')


@login_required
@require_http_methods(["POST"])
def chat_room_leave(request, room_id):
    """Rời khỏi phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    membership = ChatRoomMember.objects.filter(room=room, user=request.user).first()
    if not membership:
        messages.error(request, 'Bạn không phải thành viên của phòng này!')
        return redirect('chat_rooms_list')
    
    # Creator cannot leave their own room
    if room.created_by == request.user:
        messages.error(request, 'Bạn không thể rời khỏi phòng do chính mình tạo!')
        return redirect('chat_room_detail', room_id=room.id)
    
    membership.delete()
    
    # Create system message
    ChatMessage.objects.create(
        room=room,
        user=request.user,
        message=f'{request.user.get_full_name() or request.user.username} đã rời khỏi phòng',
        message_type='system'
    )
    
    messages.success(request, f'Đã rời khỏi phòng "{room.name}"!')
    return redirect('chat_rooms_list')


import os
from PIL import Image
from django.core.files.uploadedfile import InMemoryUploadedFile
from django.db.models import Q
import mimetypes
from django.conf import settings

# Cập nhật view chat_send_message để hỗ trợ file upload
# Thêm import này vào đầu file views.py
import time
from django.utils import timezone
from django.http import HttpResponse
import mimetypes

# Sửa lại hàm chat_send_message
@login_required
@require_http_methods(["POST"])
@csrf_exempt
def chat_send_message(request, room_id):
    """Gửi tin nhắn (text, image, file, hoặc document share)"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check membership
    membership = ChatRoomMember.objects.filter(room=room, user=request.user).first()
    if not membership:
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    if membership.is_muted:
        return JsonResponse({'error': 'Bạn đã bị cấm gửi tin nhắn!'}, status=403)
    
    # Xử lý khác nhau cho form data và JSON data
    if request.content_type and 'multipart/form-data' in request.content_type:
        # File upload
        message_text = request.POST.get('message', '').strip()
        reply_to_id = request.POST.get('reply_to')
        uploaded_file = request.FILES.get('file')
        document_id = request.POST.get('document_id')  # Share document
        
    else:
        # JSON data cho text message
        data = json.loads(request.body)
        message_text = data.get('message', '').strip()
        reply_to_id = data.get('reply_to')
        uploaded_file = None
        document_id = data.get('document_id')
    
    # Validate input
    if not message_text and not uploaded_file and not document_id:
        return JsonResponse({'error': 'Tin nhắn không thể trống!'}, status=400)
    
    reply_to = None
    if reply_to_id:
        reply_to = ChatMessage.objects.filter(id=reply_to_id, room=room).first()
    
    # Xử lý document sharing
    if document_id:
        try:
            shared_document = Document.objects.get(id=document_id, status='approved')
            if not shared_document.is_public:
                return JsonResponse({'error': 'Tài liệu này không công khai!'}, status=403)
                
            message = ChatMessage.objects.create(
                room=room,
                user=request.user,
                message=message_text or f"Chia sẻ tài liệu: {shared_document.title}",
                message_type='document_share',
                shared_document=shared_document,
                reply_to=reply_to
            )
            
        except Document.DoesNotExist:
            return JsonResponse({'error': 'Không tìm thấy tài liệu!'}, status=404)
    
    # Xử lý file upload
    elif uploaded_file:
        # Validate file size (max 50MB)
        if uploaded_file.size > 50 * 1024 * 1024:
            return JsonResponse({'error': 'File quá lớn! Tối đa 50MB.'}, status=400)
        
        # Get file info
        file_name = uploaded_file.name
        file_size = uploaded_file.size
        file_type = file_name.split('.')[-1].lower() if '.' in file_name else ''
        
        # Determine message type
        image_extensions = ['jpg', 'jpeg', 'png', 'gif', 'webp', 'bmp']
        message_type = 'image' if file_type in image_extensions else 'file'
        
        try:
            # Upload to Cloudinary
            from cloudinary.uploader import upload
            
            upload_options = {
                'folder': 'chat_files/',
                'resource_type': 'auto',
                'public_id': f"{room_id}_{request.user.id}_{int(time.time())}",
            }
            
            # For images, get dimensions và validate
            image_width = None
            image_height = None
            if message_type == 'image':
                try:
                    from PIL import Image as PILImage
                    
                    # Validate image file trước khi xử lý
                    uploaded_file.seek(0)  # Reset file pointer
                    img = PILImage.open(uploaded_file)
                    img.verify()  # Verify image integrity
                    
                    # Reopen file sau khi verify
                    uploaded_file.seek(0)
                    img = PILImage.open(uploaded_file)
                    
                    # Convert RGBA to RGB if needed
                    if img.mode in ('RGBA', 'LA', 'P'):
                        background = PILImage.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                        img = background
                    
                    image_width, image_height = img.size
                    
                    # Resize if too large
                    if image_width > 1920 or image_height > 1920:
                        img.thumbnail((1920, 1920), PILImage.Resampling.LANCZOS)
                        image_width, image_height = img.size
                        
                        # Save resized image to memory
                        from io import BytesIO
                        output = BytesIO()
                        img.save(output, format='JPEG', quality=85)
                        output.seek(0)
                        
                        uploaded_file = InMemoryUploadedFile(
                            output, None, f"{file_name.split('.')[0]}.jpg", 'image/jpeg',
                            output.getbuffer().nbytes, None
                        )
                    else:
                        # Reset file pointer for upload
                        uploaded_file.seek(0)
                        
                except Exception as e:
                    print(f"Image processing error: {e}")
                    return JsonResponse({'error': f'File ảnh không hợp lệ: {str(e)}'}, status=400)
            
            # Upload to Cloudinary
            result = upload(uploaded_file, **upload_options)
            
            message = ChatMessage.objects.create(
                room=room,
                user=request.user,
                message=message_text,
                message_type=message_type,
                file_url=result['secure_url'],
                file_name=file_name,
                file_size=file_size,
                file_type=file_type,
                image_width=image_width,
                image_height=image_height,
                reply_to=reply_to
            )
            
        except Exception as e:
            print(f"File upload error: {e}")
            return JsonResponse({'error': f'Lỗi tải file lên: {str(e)}'}, status=500)
            
    else:
        # Text message
        message = ChatMessage.objects.create(
            room=room,
            user=request.user,
            message=message_text,
            message_type='text',
            reply_to=reply_to
        )
    
    # IMPORTANT: Refresh the message from the DB to correctly load CloudinaryField properties
    message.refresh_from_db()
    
    # Update last seen
    membership.last_seen = timezone.now()
    membership.save()
    
    # Prepare response data using the centralized formatter
    message_data = format_message_data(message, request.user)
    response_data = {
        'success': True,
        'message': message_data
    }
    
    # Log activity for file uploads
    if uploaded_file:
        UserActivity.objects.create(
            user=request.user,
            action='upload_file_chat',
            description=f'Upload file "{file_name}" vào phòng "{room.name}"',
            chat_room=room
        )
    
    # Log document share activity
    if document_id:
        UserActivity.objects.create(
            user=request.user,
            action='share_document_chat', 
            description=f'Chia sẻ tài liệu "{shared_document.title}" vào phòng "{room.name}"',
            document=shared_document,
            chat_room=room
        )
    return JsonResponse(response_data)
# Thêm các views này vào file views.py của bạn

# Sửa lại view chat_room_files trong views.py

@login_required
def chat_room_files(request, room_id):
    """API lấy danh sách files đã upload trong phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check membership
    if not ChatRoomMember.objects.filter(room=room, user=request.user).exists():
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    # Get all file messages in the room
    file_messages = ChatMessage.objects.filter(
        room=room,
        message_type__in=['image', 'file'],
        is_deleted=False,
        file_url__isnull=False
    ).select_related('user').order_by('-created_at')[:50]  # Limit to latest 50 files
    
    files_data = []
    for msg in file_messages:
        # Safely handle avatar URL
        avatar_url = None
        if hasattr(msg.user, 'avatar') and msg.user.avatar:
            try:
                avatar_url = msg.user.avatar.url
            except:
                avatar_url = None
        
        files_data.append({
            'id': msg.id,
            'message_type': msg.message_type,
            'file_url': msg.file_url,
            'file_name': msg.file_name or 'Unknown',
            'file_size': msg.get_file_size_display() if hasattr(msg, 'get_file_size_display') and msg.file_size else '',
            'file_type': msg.file_type or '',
            'file_icon': msg.get_file_icon() if hasattr(msg, 'get_file_icon') else 'fa-file',
            'user': msg.user.get_full_name() or msg.user.username,
            'user_avatar': avatar_url,
            'created_at': msg.created_at.isoformat(),
            'image_width': msg.image_width,
            'image_height': msg.image_height,
        })
    
    return JsonResponse({'files': files_data})


@login_required
def chat_room_shared_documents(request, room_id):
    """API lấy danh sách tài liệu đã chia sẻ trong phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check membership
    if not ChatRoomMember.objects.filter(room=room, user=request.user).exists():
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    # Get all document share messages in the room
    document_messages = ChatMessage.objects.filter(
        room=room,
        message_type='document_share',
        is_deleted=False,
        shared_document__isnull=False
    ).select_related(
        'user', 'shared_document__university', 'shared_document__course'
    ).order_by('-created_at')[:30]  # Limit to latest 30 documents
    
    documents_data = []
    for msg in document_messages:
        doc = msg.shared_document
        
        # Safely handle avatar URL
        avatar_url = None
        if hasattr(msg.user, 'avatar') and msg.user.avatar:
            try:
                avatar_url = msg.user.avatar.url
            except:
                avatar_url = None
        
        documents_data.append({
            'id': doc.id,
            'title': doc.title,
            'description': doc.description or '',
            'university': doc.university.name if doc.university else '',
            'course': doc.course.name if doc.course else '',
            'file_type': doc.file_type or 'pdf',
            'user': msg.user.get_full_name() or msg.user.username,
            'user_avatar': avatar_url,
            'created_at': msg.created_at.isoformat(),
            'view_count': doc.view_count,
            'like_count': doc.like_count,
            'view_url': f'/documents/{doc.id}/view/',
        })
    
    return JsonResponse({'documents': documents_data})


@login_required 
def chat_room_statistics(request, room_id):
    """API lấy thống kê phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check membership
    if not ChatRoomMember.objects.filter(room=room, user=request.user).exists():
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    # Get statistics
    total_messages = ChatMessage.objects.filter(room=room, is_deleted=False).count()
    total_files = ChatMessage.objects.filter(
        room=room, 
        message_type__in=['image', 'file'], 
        is_deleted=False
    ).count()
    total_documents = ChatMessage.objects.filter(
        room=room, 
        message_type='document_share', 
        is_deleted=False
    ).count()
    
    # Most active users
    from django.db.models import Count
    active_users = ChatMessage.objects.filter(
        room=room, 
        is_deleted=False
    ).values(
        'user__username', 
        'user__first_name', 
        'user__last_name'
    ).annotate(
        message_count=Count('id')
    ).order_by('-message_count')[:5]
    
    # Recent activity by day
    from django.utils import timezone
    from datetime import timedelta
    
    last_7_days = timezone.now() - timedelta(days=7)
    daily_activity = ChatMessage.objects.filter(
        room=room,
        is_deleted=False,
        created_at__gte=last_7_days
    ).extra({
        'day': 'date(created_at)'
    }).values('day').annotate(
        count=Count('id')
    ).order_by('day')
    
    statistics = {
        'total_messages': total_messages,
        'total_files': total_files,
        'total_documents': total_documents,
        'active_users': list(active_users),
        'daily_activity': list(daily_activity)
    }
    
    return JsonResponse(statistics)
# Thêm view mới để handle download file


# Thêm helper function này vào đầu views.py hoặc tạo file utils.py riêng

def get_safe_cloudinary_url(cloudinary_field):
    """
    Safely get URL from Cloudinary field
    Returns None if field is empty or invalid
    """
    if not cloudinary_field:
        return None
    
    try:
        # Try to get URL - this might fail if field is empty or corrupted
        return cloudinary_field.url
    except (AttributeError, ValueError, TypeError):
        return None


def serialize_message_for_json(message):
    """
    Safely serialize a ChatMessage object for JSON response
    """
    avatar_url = get_safe_cloudinary_url(message.user.avatar) if hasattr(message.user, 'avatar') else None
    
    return {
        'id': message.id,
        'message_type': message.message_type,
        'file_url': get_safe_cloudinary_url(message.file_url) or '',
        'file_name': message.file_name or 'Unknown',
        'file_size': message.get_file_size_display() if hasattr(message, 'get_file_size_display') and message.file_size else '',
        'file_type': message.file_type or '',
        'file_icon': message.get_file_icon() if hasattr(message, 'get_file_icon') else 'fa-file',
        'user': message.user.get_full_name() or message.user.username,
        'user_avatar': avatar_url,
        'created_at': message.created_at.isoformat(),
        'image_width': message.image_width,
        'image_height': message.image_height,
    }

# Sau đó sử dụng helper function này trong view:

@login_required
def chat_room_files(request, room_id):
    """API lấy danh sách files đã upload trong phòng chat"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check membership
    if not ChatRoomMember.objects.filter(room=room, user=request.user).exists():
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    try:
        # Get all file messages in the room
        file_messages = ChatMessage.objects.filter(
            room=room,
            message_type__in=['image', 'file'],
            is_deleted=False,
            file_url__isnull=False
        ).select_related('user').order_by('-created_at')[:50]
        
        files_data = [serialize_message_for_json(msg) for msg in file_messages]
        
        return JsonResponse({'files': files_data})
        
    except Exception as e:
        import logging
        logging.error(f"Error in chat_room_files: {str(e)}")
        return JsonResponse({'error': 'Lỗi server', 'files': []}, status=500)
@login_required
def chat_file_download(request, room_id, message_id):
    """Download file from chat message by proxying it from Cloudinary."""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check membership
    if not ChatRoomMember.objects.filter(room=room, user=request.user).exists():
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    message = get_object_or_404(ChatMessage, id=message_id, room=room, message_type__in=['file', 'image'])
    
    if not message.file_url:
        return JsonResponse({'error': 'File không tồn tại!'}, status=404)

    try:
        import cloudinary.utils
        public_id = message.file_url.public_id
        
        # Determine the correct resource type based on the message type
        resource_type = 'raw' if message.message_type == 'file' else 'image'
        
        # For 'raw' files (non-images), the file extension must be specified
        # in the delivery URL. We use the 'format' parameter for this.
        file_format = message.file_type if resource_type == 'raw' else None

        # Build the correct delivery URL
        direct_url = cloudinary.utils.cloudinary_url(
            public_id,
            resource_type=resource_type,
            format=file_format
        )[0]
        
        # Make a streaming request to the file URL
        cloudinary_response = requests.get(direct_url, stream=True)
        
        # Check if the request to Cloudinary was successful
        cloudinary_response.raise_for_status()

        # Get content type from Cloudinary's response
        content_type = cloudinary_response.headers.get('Content-Type', 'application/octet-stream')

        # Create a streaming response to send to the user
        response = StreamingHttpResponse(
            cloudinary_response.iter_content(chunk_size=8192),
            content_type=content_type
        )
        
        # URL-encode the filename for the Content-Disposition header for safety
        safe_filename = quote(message.file_name)
        response['Content-Disposition'] = f"attachment; filename*=UTF-8''{safe_filename}"
        response['Content-Length'] = cloudinary_response.headers.get('Content-Length')

        return response

    except requests.exceptions.RequestException as e:
        # If there's an error fetching from Cloudinary (e.g., 404)
        return HttpResponse(f"Error fetching file from storage: {e}", status=502) # 502 Bad Gateway
    except (AttributeError, TypeError):
        # Fallback for old data where file_url might be a full URL, attempt to redirect
        return redirect(str(message.file_url))
# API tìm kiếm tài liệu cho chat
def normalize_vietnamese(text):
    """Normalize Vietnamese text for fuzzy search"""
    return unidecode(text.lower()).strip()

@login_required
def chat_search_documents(request, room_id):
    """API to search for documents to share in chat, with fuzzy search and initial listing."""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)

    if not ChatRoomMember.objects.filter(room=room, user=request.user).exists():
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)

    query = request.GET.get('q', '').strip()
    
    base_query = Document.objects.filter(status='approved', is_public=True)
    
    # Filter by room's university/course if available
    if room.university:
        base_query = base_query.filter(university=room.university)
    if room.course:
        base_query = base_query.filter(course=room.course)

    if not query:
        # If query is empty, return the latest documents
        documents = base_query.order_by('-created_at')[:20]
    else:
        # If there is a query, perform fuzzy search
        query_norm = normalize_vietnamese(query)
        
        # Fetch all potential documents to perform fuzzy matching in Python
        all_documents = base_query.select_related('university', 'course', 'uploaded_by')
        
        matched_docs = []
        for doc in all_documents:
            title_norm = normalize_vietnamese(doc.title)
            # Calculate similarity ratio
            score = ratio(query_norm, title_norm)
            
            # Also check for simple containment for partial matches
            if query_norm in title_norm:
                score = max(score, 0.8) # Boost score for containment

            if score >= 0.6:  # Threshold for matching
                matched_docs.append({
                    'doc': doc,
                    'score': score
                })
        
        # Sort by score in descending order
        matched_docs.sort(key=lambda x: x['score'], reverse=True)
        
        # Get the final list of document objects
        documents = [item['doc'] for item in matched_docs[:20]]

        # Save search history only for non-empty queries
        ChatDocumentSearch.objects.create(
            room=room,
            user=request.user,
            query=query,
            results_count=len(documents)
        )

    documents_data = []
    for doc in documents:
        documents_data.append({
            'id': doc.id,
            'title': doc.title,
            'description': doc.description or '',
            'university': doc.university.name if doc.university else '',
            'course': doc.course.name if doc.course else '',
            'uploaded_by': doc.uploaded_by.get_full_name() or doc.uploaded_by.username,
            'file_type': doc.file_type or 'pdf',
            'view_count': doc.view_count,
            'like_count': doc.like_count,
            'created_at': doc.created_at.strftime('%d/%m/%Y'),
            'view_url': f'/documents/{doc.id}/view/',
        })
    
    return JsonResponse({'documents': documents_data})


@login_required
def chat_load_messages(request, room_id):
    """Load tin nhắn - hỗ trợ cả polling (tin nhắn mới) và pagination (tin nhắn cũ)"""
    try:
        room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
        
        # Check membership
        membership = ChatRoomMember.objects.filter(room=room, user=request.user).first()
        if not membership:
            return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
        
        # Update last seen
        membership.last_seen = timezone.now()
        membership.save()
        
        # Xử lý polling - lấy tin nhắn mới từ last_id
        last_id = request.GET.get('last_id')
        if last_id:
            try:
                last_id = int(last_id)
                new_messages = ChatMessage.objects.filter(
                    room=room,
                    is_deleted=False,
                    id__gt=last_id
                ).select_related(
                    'user', 'reply_to__user', 'shared_document__university', 'shared_document__course'
                ).order_by('created_at')[:20]
                
                messages_data = []
                for msg in new_messages:
                    message_data = format_message_data(msg, request.user)
                    messages_data.append(message_data)
                
                return JsonResponse({'messages': messages_data})
                
            except (ValueError, TypeError):
                return JsonResponse({'error': 'Invalid last_id parameter'}, status=400)
        
        # Xử lý pagination - load tin nhắn cũ
        offset = int(request.GET.get('offset', 0))
        limit = 20
        
        messages_list = ChatMessage.objects.filter(
            room=room,
            is_deleted=False
        ).select_related(
            'user', 'reply_to__user', 'shared_document__university', 'shared_document__course'
        ).order_by('-created_at')[offset:offset+limit]
        
        messages_data = []
        for msg in reversed(messages_list):
            message_data = format_message_data(msg, request.user)
            messages_data.append(message_data)
        
        # Check if there are more messages
        total_messages = ChatMessage.objects.filter(room=room, is_deleted=False).count()
        has_more = total_messages > offset + limit
        
        return JsonResponse({
            'messages': messages_data,
            'has_more': has_more,
            'total': total_messages
        })
        
    except Exception as e:
        print(f"chat_load_messages error: {str(e)}")
        import traceback
        traceback.print_exc()
        return JsonResponse({'error': 'Lỗi server khi tải tin nhắn'}, status=500)


def format_message_data(msg, current_user):
    """Helper function để format message data"""
    message_data = {
        'id': msg.id,
        'user': msg.user.get_full_name() or msg.user.username,
        'user_avatar': get_safe_cloudinary_url(msg.user.avatar),
        'message': msg.message,
        'message_type': msg.message_type,
        'created_at': msg.created_at.strftime('%H:%M'),
        'is_own': msg.user == current_user,
        'is_edited': msg.is_edited,
        'reply_to': None
    }
    
    # Add reply_to data if exists
    if msg.reply_to:
        reply_message = msg.reply_to.message or ''
        if len(reply_message) > 50:
            reply_message = reply_message[:50] + '...'
        
        message_data['reply_to'] = {
            'user': msg.reply_to.user.get_full_name() or msg.reply_to.user.username,
            'message': reply_message
        }
    
    # Add specific data based on message type
    if msg.message_type == 'image':
        message_data.update({
            'file_url': get_safe_cloudinary_url(msg.file_url),
            'file_name': msg.file_name,
            'image_width': msg.image_width,
            'image_height': msg.image_height,
        })
    elif msg.message_type == 'file':
        message_data.update({
            'file_url': get_safe_cloudinary_url(msg.file_url),
            'file_name': msg.file_name,
            'file_size': msg.get_file_size_display(),
            'file_icon': msg.get_file_icon(),
        })
    elif msg.message_type == 'document_share' and msg.shared_document:
        message_data.update({
            'document': {
                'id': msg.shared_document.id,
                'title': msg.shared_document.title,
                'description': msg.shared_document.description or '',
                'university': msg.shared_document.university.name if msg.shared_document.university else '',
                'course': msg.shared_document.course.name if msg.shared_document.course else '',
                'file_type': msg.shared_document.file_type or 'pdf',
                'view_url': f'/documents/{msg.shared_document.id}/view/',
            }
        })
    
    return message_data

@login_required
def chat_room_members(request, room_id):
    """API lấy danh sách thành viên"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check membership
    if not ChatRoomMember.objects.filter(room=room, user=request.user).exists():
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    members = ChatRoomMember.objects.filter(room=room).select_related('user').order_by('-role', '-last_seen')
    
    members_data = []
    for member in members:
        members_data.append({
            'id': member.user.id,
            'username': member.user.username,
            'full_name': member.user.get_full_name(),
            'avatar': member.user.avatar.url if member.user.avatar else None,
            'role': member.role,
            'role_display': member.get_role_display(),
            'joined_at': member.joined_at.strftime('%d/%m/%Y'),
            'last_seen': member.last_seen.strftime('%H:%M %d/%m') if member.last_seen else None,
            'is_online': member.last_seen and (timezone.now() - member.last_seen).seconds < 300,
            'is_muted': member.is_muted,
        })
    
    return JsonResponse({'members': members_data})


@login_required
@require_http_methods(["POST"])
def chat_room_invite(request, room_id):
    """Mời người khác vào phòng"""
    room = get_object_or_404(ChatRoom, id=room_id, is_active=True)
    
    # Check permissions
    membership = ChatRoomMember.objects.filter(room=room, user=request.user).first()
    if not membership:
        return JsonResponse({'error': 'Bạn không phải thành viên của phòng này!'}, status=403)
    
    username = request.POST.get('username', '').strip()
    if not username:
        return JsonResponse({'error': 'Vui lòng nhập tên người dùng!'}, status=400)
    
    try:
        invited_user = User.objects.get(username=username)
    except User.DoesNotExist:
        return JsonResponse({'error': 'Không tìm thấy người dùng!'}, status=404)
    
    # Check if already member
    if ChatRoomMember.objects.filter(room=room, user=invited_user).exists():
        return JsonResponse({'error': 'Người dùng đã là thành viên của phòng!'}, status=400)
    
    # Check room capacity
    members_count = ChatRoomMember.objects.filter(room=room).count()
    if members_count >= room.max_members:
        return JsonResponse({'error': 'Phòng đã đầy!'}, status=400)
    
    # Add member
    ChatRoomMember.objects.create(
        room=room,
        user=invited_user,
        role='member'
    )
    
    # Create system message
    ChatMessage.objects.create(
        room=room,
        user=request.user,
        message=f'{request.user.get_full_name() or request.user.username} đã mời {invited_user.get_full_name() or invited_user.username} vào phòng',
        message_type='system'
    )
    
    # Create notification
    Notification.objects.create(
        user=invited_user,
        title='Lời mời tham gia phòng chat',
        message=f'Bạn được mời tham gia phòng "{room.name}" bởi {request.user.get_full_name() or request.user.username}',
        notification_type='info',
        chat_room=room
    )
    
    return JsonResponse({'success': True, 'message': 'Mời thành công!'})

import requests
import json
import base64
import time
import uuid
import traceback
from django.conf import settings
from google.oauth2 import id_token
from google.auth.transport import requests as google_requests
from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods
from django.core.files.base import ContentFile
from django.contrib import messages
from PIL import Image
import io
import tempfile
import os

# Import thêm các thư viện xử lý file
import PyPDF2
import docx
from pptx import Presentation
import pandas as pd
import openpyxl

# Import models
from .models import (
    AIImageSolution, AIConversation, AIConversationMessage, AIImageSolutionLike,
    Document, ChatRoom, Course, University, User, ChatRoomMember
)
from django.db.models import Q, Count
from django.contrib.postgres.search import SearchVector


def verify_google_token(token):
    """Verify Google ID token and return user info"""
    try:
        # Verify token with Google
        idinfo = id_token.verify_oauth2_token(
            token, 
            google_requests.Request(), 
            settings.GOOGLE_CLIENT_ID
        )
        
        if idinfo['iss'] not in ['accounts.google.com', 'https://accounts.google.com']:
            raise ValueError('Wrong issuer.')
            
        return idinfo
    except ValueError as e:
        print(f"Token verification failed: {e}")
        return None

# Gemini API configuration
GEMINI_API_KEY = "AIzaSyAK_FEeBSpDQq5P6THmj1fp89jRPNVPX7E"
GEMINI_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"
# Supported file types
SUPPORTED_FILE_TYPES = {
    'application/pdf': '.pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/msword': '.doc',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
    'application/vnd.ms-powerpoint': '.ppt',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
    'application/vnd.ms-excel': '.xls',
    'text/plain': '.txt',
    'text/csv': '.csv'
}

def search_documents_for_ai(query, user, limit=5):
    """Search documents relevant to user query with debug and fallback"""
    try:
        print(f"=== DEBUG SEARCH DOCUMENTS ===")
        print(f"Query: '{query}'")
        print(f"User: {user.username}")
        
        # Đếm tổng số documents trong DB
        total_docs = Document.objects.count()
        print(f"Total documents in DB: {total_docs}")
        
        # Đếm documents approved và public
        approved_public = Document.objects.filter(
            status='approved', 
            is_public=True
        ).count()
        print(f"Approved + Public documents: {approved_public}")
        
        # Test search với điều kiện đơn giản trước
        simple_search = Document.objects.filter(
            status='approved',
            is_public=True,
            title__icontains=query
        )
        print(f"Simple title search results: {simple_search.count()}")
        
        # Nếu không có kết quả với approved/public, thử search tất cả
        if simple_search.count() == 0:
            all_search = Document.objects.filter(
                title__icontains=query
            )
            print(f"Search all documents (ignore status): {all_search.count()}")
            
            # In ra một vài documents để kiểm tra
            sample_docs = Document.objects.all()[:3]
            for doc in sample_docs:
                print(f"Sample doc: {doc.title}, status: {doc.status}, public: {doc.is_public}")
        
        # Search chính với điều kiện đầy đủ
        documents = Document.objects.filter(
            Q(status='approved') & Q(is_public=True) &
            (Q(title__icontains=query) |
             Q(description__icontains=query) |
             Q(ai_summary__icontains=query) |
             Q(ai_keywords__icontains=query))
        ).select_related('course', 'university', 'uploaded_by').order_by('-view_count', '-created_at')[:limit]
        
        print(f"Final search results: {documents.count()}")
        
        # FIXED: Fallback nếu không tìm thấy
        if documents.count() == 0:
            print("No results with query, trying fallback...")
            # Fallback: lấy tất cả approved docs
            documents = Document.objects.filter(
                status='approved',
                is_public=True
            ).select_related('course', 'university', 'uploaded_by')[:limit]
            print(f"Fallback results: {documents.count()}")
        
        # FIXED: Debug approved documents
        if documents.count() > 0:
            print("=== APPROVED DOCUMENTS FOUND ===")
            for doc in documents:
                print(f"- ID: {doc.id}, Title: '{doc.title}', Course: {doc.course}, University: {doc.university}")
        
        # Format results for AI
        results = []
        for doc in documents:
            try:
                result_item = {
                    'id': doc.id,
                    'title': doc.title,
                    'description': doc.description[:200] if doc.description else '',
                    'course': f"{doc.course.code} - {doc.course.name}" if doc.course else 'No course',
                    'university': doc.university.name if doc.university else 'No university',
                    'document_type': doc.get_document_type_display() if hasattr(doc, 'get_document_type_display') else 'Document',
                    'view_count': getattr(doc, 'view_count', 0),
                    'like_count': getattr(doc, 'like_count', 0),
                    'url': f"/documents/{doc.id}/view/"
                }
                results.append(result_item)
                print(f"Formatted doc: {result_item['title']}")
            except Exception as e:
                print(f"Error formatting document {doc.id}: {e}")
                continue
            
        print(f"Formatted results: {len(results)}")
        print("=== END DEBUG ===")
        
        return results
    except Exception as e:
        print(f"Error searching documents: {e}")
        import traceback
        traceback.print_exc()
        return []
def search_chat_rooms_for_ai(query, user, limit=5):
    """Search chat rooms relevant to user query with debug"""
    try:
        print(f"=== DEBUG SEARCH CHAT ROOMS ===")
        print(f"Query: '{query}'")
        print(f"User: {user.username}")
        
        # Đếm tổng số chat rooms
        total_rooms = ChatRoom.objects.count()
        print(f"Total chat rooms in DB: {total_rooms}")
        
        # Đếm active rooms
        active_rooms = ChatRoom.objects.filter(
            is_active=True,
            room_type__in=['public', 'group']
        ).count()
        print(f"Active public/group rooms: {active_rooms}")
        
        # Test search đơn giản
        simple_search = ChatRoom.objects.filter(
            is_active=True,
            room_type__in=['public', 'group'],
            name__icontains=query
        )
        print(f"Simple name search results: {simple_search.count()}")
        
        # Nếu không có, thử search tất cả
        if simple_search.count() == 0:
            all_search = ChatRoom.objects.filter(
                name__icontains=query
            )
            print(f"Search all rooms (ignore conditions): {all_search.count()}")
            
            # Sample rooms
            sample_rooms = ChatRoom.objects.all()[:3]
            for room in sample_rooms:
                print(f"Sample room: {room.name}, active: {room.is_active}, type: {room.room_type}")
        
        # Search chính - FIXED: đổi 'chatroomMember' thành 'chatroommember'
        rooms = ChatRoom.objects.filter(
            Q(is_active=True) & Q(room_type__in=['public', 'group']) &
            (Q(name__icontains=query) |
             Q(description__icontains=query) |
             Q(course__name__icontains=query) |
             Q(course__code__icontains=query) |
             Q(university__name__icontains=query))
        ).select_related('course', 'university', 'created_by').annotate(
            member_count=Count('chatroommember')  # FIXED: lowercase
        ).order_by('-member_count', '-created_at')[:limit]
        
        print(f"Final search results: {rooms.count()}")
        
        # FIXED: Fallback nếu không tìm thấy
        if rooms.count() == 0:
            print("No results with query, trying fallback...")
            rooms = ChatRoom.objects.filter(
                is_active=True,
                room_type__in=['public', 'group']
            ).select_related('course', 'university', 'created_by').annotate(
                member_count=Count('chatroommember')
            )[:limit]
            print(f"Fallback results: {rooms.count()}")
        
        # Format results for AI
        results = []
        for room in rooms:
            try:
                result_item = {
                    'id': room.id,
                    'name': room.name,
                    'description': room.description[:200] if room.description else '',
                    'course': f"{room.course.code} - {room.course.name}" if room.course else 'Chung',
                    'university': room.university.name if room.university else 'Chung',
                    'room_type': room.get_room_type_display() if hasattr(room, 'get_room_type_display') else room.room_type,
                    'member_count': getattr(room, 'member_count', 0),
                    'url': f"/chat/room/{room.id}/"
                }
                results.append(result_item)
                print(f"Formatted room: {result_item['name']}")
            except Exception as e:
                print(f"Error formatting room {room.id}: {e}")
                continue
            
        print(f"Formatted results: {len(results)}")
        print("=== END DEBUG ===")
        
        return results
    except Exception as e:
        print(f"Error searching chat rooms: {e}")
        import traceback
        traceback.print_exc()
        return []

def search_courses_for_ai(query, user, limit=5):
    """Search coding courses relevant to user query"""
    try:
        from .models import CodeCourse
        
        print(f"=== DEBUG SEARCH COURSES FOR AI ===")
        print(f"Query: '{query}'")
        print(f"User: {user.username}")
        
        # Count total courses
        total_courses = CodeCourse.objects.filter(status='published').count()
        print(f"Total published courses: {total_courses}")
        
        # Search courses
        courses = CodeCourse.objects.filter(
            Q(status='published') &
            (Q(title__icontains=query) |
             Q(description__icontains=query) |
             Q(language__display_name__icontains=query) |
             Q(language__name__icontains=query) |
             Q(difficulty__icontains=query))
        ).select_related('language', 'created_by', 'university').order_by('-enrollment_count', '-rating_average')[:limit]
        
        print(f"Search results: {courses.count()}")
        
        # Fallback if no results
        if courses.count() == 0:
            print("No results with query, trying fallback...")
            courses = CodeCourse.objects.filter(
                status='published'
            ).select_related('language', 'created_by', 'university').order_by('-enrollment_count', '-rating_average')[:limit]
            print(f"Fallback results: {courses.count()}")
        
        # Format results
        results = []
        for course in courses:
            try:
                result_item = {
                    'id': course.id,
                    'title': course.title,
                    'description': course.description[:200] if course.description else '',
                    'language': course.language.display_name,
                    'difficulty': course.get_difficulty_display() if hasattr(course, 'get_difficulty_display') else course.difficulty,
                    'enrollment_count': course.enrollment_count,
                    'rating': float(course.rating_average) if course.rating_average else 0,
                    'url': f"/code/courses/{course.slug}/"
                }
                results.append(result_item)
                print(f"Formatted course: {result_item['title']}")
            except Exception as e:
                print(f"Error formatting course {course.id}: {e}")
                continue
        
        print(f"Formatted results: {len(results)}")
        print("=== END DEBUG ===")
        
        return results
    except Exception as e:
        print(f"Error searching courses: {e}")
        import traceback
        traceback.print_exc()
        return []
def search_documents_simple(query, user, limit=5):
    """Simple search function as backup"""
    try:
        documents = Document.objects.filter(
            Q(title__icontains=query) |
            Q(description__icontains=query)
        ).select_related('course', 'university', 'uploaded_by')[:limit]
        
        results = []
        for doc in documents:
            results.append({
                'id': doc.id,
                'title': doc.title,
                'description': doc.description[:200] if doc.description else '',
                'course': f"{doc.course.code} - {doc.course.name}" if doc.course else 'No course',
                'university': doc.university.name if doc.university else 'No university',
                'url': f"/documents/{doc.id}/view/"
            })
        return results
    except Exception as e:
        print(f"Simple search error: {e}")
        return []

def search_chat_rooms_simple(query, user, limit=5):
    """Simple search function as backup"""
    try:
        rooms = ChatRoom.objects.filter(
            Q(name__icontains=query) |
            Q(description__icontains=query)
        ).select_related('course', 'university', 'created_by')[:limit]
        
        results = []
        for room in rooms:
            results.append({
                'id': room.id,
                'name': room.name,
                'description': room.description[:200] if room.description else '',
                'course': f"{room.course.code} - {room.course.name}" if room.course else 'Chung',
                'university': room.university.name if room.university else 'Chung',
                'url': f"/chat/room/{room.id}/"
            })
        return results
    except Exception as e:
        print(f"Simple search error: {e}")
        return []
def get_user_courses_and_interests(user):
    """Get user's courses and interests for better recommendations"""
    try:
        # Get courses from user's documents
        user_courses = Course.objects.filter(
            document__uploaded_by=user
        ).distinct().values_list('name', 'code')
        
        # Get courses from user's chat rooms
        user_chat_courses = Course.objects.filter(
            chatroom__chatroomMember__user=user
        ).distinct().values_list('name', 'code')
        
        # Combine and format
        all_courses = list(user_courses) + list(user_chat_courses)
        course_info = [f"{code} - {name}" for code, name in set(all_courses)]
        
        return course_info
    except Exception as e:
        print(f"Error getting user courses: {e}")
        return []

def enhance_ai_prompt_with_context(user_message, user, conversation_type='text'):
    """Enhance AI prompt with database context"""
    try:
        # Get user context
        user_courses = get_user_courses_and_interests(user)
        
        # Search for relevant documents, chat rooms, and courses
        docs = search_documents_for_ai(user_message, user, limit=3)
        chat_rooms = search_chat_rooms_for_ai(user_message, user, limit=3)
        courses = search_courses_for_ai(user_message, user, limit=3)
        
        # Build enhanced context
        context_parts = []
        
        if user_courses:
            context_parts.append(f"Thông tin người dùng: Đang học/quan tâm các môn: {', '.join(user_courses[:5])}")
        
        if docs:
            doc_info = []
            for doc in docs:
                doc_info.append(f"- [{doc['title']}]({doc['url']}) - {doc['course']} ({doc['view_count']} lượt xem)")
            context_parts.append(f"Tài liệu liên quan:\n" + "\n".join(doc_info))
        
        if chat_rooms:
            room_info = []
            for room in chat_rooms:
                room_info.append(f"- [{room['name']}]({room['url']}) - {room['course']} ({room['member_count']} thành viên)")
            context_parts.append(f"Phòng chat liên quan:\n" + "\n".join(room_info))
        
        if courses:
            course_info = []
            for course in courses:
                course_info.append(f"- [{course['title']}]({course['url']}) - {course['language']}, {course['difficulty']} ({course['enrollment_count']} học viên, ⭐{course['rating']:.1f})")
            context_parts.append(f"Khóa học lập trình:\n" + "\n".join(course_info))
        
        return "\n\n".join(context_parts) if context_parts else ""
        
    except Exception as e:
        print(f"Error enhancing prompt: {e}")
        return ""
def call_gemini_api_enhanced(messages, image_data=None, user=None):
    """Enhanced Gemini API call with database context"""
    try:
        headers = {
            'Content-Type': 'application/json',
        }
        
        # Prepare contents for API
        contents = []
        
        # FIXED: Enhanced system prompt - chỉ gợi ý tài liệu có thật
        system_prompt = '''Bạn là AI assistant thông minh của một website chia sẻ tài liệu học tập và khóa học lập trình.
        
        Khả năng của bạn:
        1. Giải thích kiến thức học tập (toán, lý, hóa, văn, anh, v.v.)
        2. Giải bài tập và hướng dẫn từng bước
        3. Trả lời câu hỏi thường thức
        4. Gợi ý tài liệu, phòng chat, và khóa học lập trình CÓ TRONG HỆ THỐNG
        5. Hỗ trợ học tập và nghiên cứu
        
        QUY TẮC QUAN TRỌNG:
        - CHỈ gợi ý tài liệu, phòng chat, và khóa học được cung cấp trong phần "Thông tin từ hệ thống"
        - KHÔNG tự tạo ra tài liệu, phòng chat, hoặc khóa học không tồn tại
        - Nếu không có tài liệu phù hợp trong hệ thống, nói rằng "Hiện tại chưa có tài liệu phù hợp trong hệ thống"
        - Nếu không có phòng chat phù hợp, nói rằng "Hiện tại chưa có phòng chat phù hợp trong hệ thống"
        - Nếu không có khóa học phù hợp, nói rằng "Hiện tại chưa có khóa học phù hợp trong hệ thống"
        
        Khi có thông tin từ hệ thống:
        📚 **Tài liệu liên quan:**
        [Chỉ liệt kê tài liệu được cung cấp với đúng link]
        
        💬 **Phòng chat để thảo luận:**
        [Chỉ liệt kê phòng chat được cung cấp với đúng link]
        
        💻 **Khóa học lập trình:**
        [Chỉ liệt kê khóa học được cung cấp với đúng link]
        
        Khi không có thông tin từ hệ thống:
        - Trả lời câu hỏi bình thường
        - Nói rõ là hiện tại chưa có tài liệu/phòng chat/khóa học phù hợp
        - KHÔNG đưa ra gợi ý tài liệu/phòng chat/khóa học giả tạo
        
        Luôn trả lời bằng tiếng Việt, thân thiện và khích lệ người học.
        '''
        
        # Add conversation history
        for i, msg in enumerate(messages):
            if msg['role'] == 'system':
                continue  # Will be handled separately
                
            content = {
                "role": "user" if msg['role'] == 'user' else "model",
                "parts": [{"text": msg['content']}]
            }
            
            # Enhance the last user message with database context
            if i == len(messages) - 1 and msg['role'] == 'user' and user:
                db_context = enhance_ai_prompt_with_context(msg['content'], user)
                if db_context:
                    enhanced_content = f"{msg['content']}\n\n--- Thông tin từ hệ thống ---\n{db_context}"
                    content['parts'] = [{"text": enhanced_content}]
                    print(f"Enhanced user message with DB context: {len(db_context)} chars")
                else:
                    # IMPORTANT: Thêm thông báo không có dữ liệu
                    enhanced_content = f"{msg['content']}\n\n--- Thông tin từ hệ thống ---\nKhông tìm thấy tài liệu hoặc phòng chat phù hợp trong hệ thống."
                    content['parts'] = [{"text": enhanced_content}]
                    print("No DB context found - added no data message")
            
            contents.append(content)
        
        # Add system prompt as the first user message
        if contents:
            contents.insert(0, {
                "role": "user",
                "parts": [{"text": system_prompt}]
            })
            contents.insert(1, {
                "role": "model", 
                "parts": [{"text": "Tôi hiểu. Tôi sẽ chỉ gợi ý tài liệu và phòng chat có thật trong hệ thống. Nếu không có, tôi sẽ nói rõ là chưa có dữ liệu phù hợp thay vì tự tạo gợi ý."}]
            })
        
        # Add image if provided
        if image_data and contents:
            # Add image to the last user message
            for content in reversed(contents):
                if content['role'] == 'user':
                    content['parts'].append({
                        "inline_data": image_data
                    })
                    break
        
        payload = {
            "contents": contents,
            "generationConfig": {
                "temperature": 0.7,
                "topK": 40,
                "topP": 0.95,
                "maxOutputTokens": 8192,
            },
            "safetySettings": [
                {
                    "category": "HARM_CATEGORY_HARASSMENT",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                },
                {
                    "category": "HARM_CATEGORY_HATE_SPEECH",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                },
                {
                    "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                },
                {
                    "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                }
            ]
        }
        
        response = requests.post(
            f"{GEMINI_URL}?key={GEMINI_API_KEY}",
            headers=headers,
            data=json.dumps(payload),
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            if 'candidates' in result and len(result['candidates']) > 0:
                content = result['candidates'][0]['content']['parts'][0]['text']
                return {
                    'success': True,
                    'content': content,
                    'usage': result.get('usageMetadata', {})
                }
            else:
                return {
                    'success': False,
                    'error': 'No response from AI'
                }
        else:
            return {
                'success': False,
                'error': f'API Error: {response.status_code} - {response.text}'
            }
            
    except Exception as e:
        return {
            'success': False,
            'error': f'Exception: {str(e)}'
        }

def enhance_ai_prompt_with_context_strict(user_message, user, conversation_type='text'):
    """Enhanced AI prompt - CHỈ trả về nếu có dữ liệu thật"""
    try:
        print(f"=== STRICT ENHANCE PROMPT ===")
        print(f"User message: '{user_message[:100]}...'")
        
        # Get user context
        user_courses = get_user_courses_and_interests(user)
        print(f"User courses: {user_courses}")
        
        # Search documents and chat rooms
        docs = search_documents_for_ai(user_message, user, limit=3)
        chat_rooms = search_chat_rooms_for_ai(user_message, user, limit=3)
        
        print(f"Search results - Docs: {len(docs)}, Rooms: {len(chat_rooms)}")
        
        # CHỈ BUILD CONTEXT NỀU CÓ DỮ LIỆU THẬT
        context_parts = []
        
        # User courses luôn có thể thêm
        if user_courses:
            context_parts.append(f"Thông tin người dùng: Đang học/quan tâm các môn: {', '.join(user_courses[:5])}")
        
        # CHỈ thêm documents nếu tìm thấy trong DB
        if docs:
            doc_info = []
            for doc in docs:
                doc_info.append(f"- [{doc['title']}]({doc['url']}) - {doc['course']}")
            context_parts.append(f"Tài liệu có trong hệ thống:\n" + "\n".join(doc_info))
        
        # CHỈ thêm chat rooms nếu tìm thấy trong DB  
        if chat_rooms:
            room_info = []
            for room in chat_rooms:
                room_info.append(f"- [{room['name']}]({room['url']}) - {room['course']}")
            context_parts.append(f"Phòng chat có trong hệ thống:\n" + "\n".join(room_info))
        
        enhanced_context = "\n\n".join(context_parts) if context_parts else ""  
        
        print(f"Final context length: {len(enhanced_context)}")
        print(f"Has docs: {bool(docs)}, Has rooms: {bool(chat_rooms)}")
        print("=== END STRICT ENHANCE ===")
        
        return enhanced_context
        
    except Exception as e:
        print(f"Error enhancing prompt: {e}")
        import traceback
        traceback.print_exc()
        return ""

def extract_text_from_file(file):
    """Extract text content from various file types"""
    try:
        # Get file extension
        file_extension = None
        if hasattr(file, 'name') and file.name:
            file_extension = os.path.splitext(file.name)[1].lower()
        
        # Reset file pointer
        file.seek(0)
        
        # PDF files
        if file_extension == '.pdf':
            return extract_pdf_text(file)
        
        # Word documents
        elif file_extension in ['.docx']:
            return extract_docx_text(file)
        
        # PowerPoint presentations
        elif file_extension in ['.pptx']:
            return extract_pptx_text(file)
        
        # Excel files
        elif file_extension in ['.xlsx', '.xls']:
            return extract_excel_text(file)
        
        # Text files
        elif file_extension in ['.txt', '.csv']:
            return extract_text_file(file)
        
        # DOC files (older format) - requires python-docx2txt or similar
        elif file_extension == '.doc':
            return "Định dạng DOC không được hỗ trợ trực tiếp. Vui lòng chuyển đổi sang DOCX."
        
        # PPT files (older format)
        elif file_extension == '.ppt':
            return "Định dạng PPT không được hỗ trợ trực tiếp. Vui lòng chuyển đổi sang PPTX."
        
        else:
            return f"Loại file {file_extension} không được hỗ trợ."
            
    except Exception as e:
        print(f"Error extracting text from file: {e}")
        return f"Lỗi khi đọc file: {str(e)}"


def extract_pdf_text(file):
    """Extract text from PDF file"""
    try:
        # Create a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            # Write file content to temporary file
            file.seek(0)
            tmp_file.write(file.read())
            tmp_file.flush()
            
            # Read PDF
            text = ""
            with open(tmp_file.name, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text += f"\n--- Trang {page_num + 1} ---\n"
                            text += page_text + "\n"
                    except Exception as e:
                        text += f"\n--- Lỗi đọc trang {page_num + 1}: {str(e)} ---\n"
            
            # Cleanup
            os.unlink(tmp_file.name)
            
            return text.strip() if text.strip() else "Không thể trích xuất text từ PDF này."
            
    except Exception as e:
        return f"Lỗi khi đọc PDF: {str(e)}"


def extract_docx_text(file):
    """Extract text from DOCX file"""
    try:
        file.seek(0)
        doc = docx.Document(file)
        
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + " | "
                text += "\n"
        
        return text.strip() if text.strip() else "Document này không chứa text có thể đọc được."
        
    except Exception as e:
        return f"Lỗi khi đọc DOCX: {str(e)}"


def extract_pptx_text(file):
    """Extract text from PPTX file"""
    try:
        file.seek(0)
        presentation = Presentation(file)
        
        text = ""
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        
        return text.strip() if text.strip() else "Presentation này không chứa text có thể đọc được."
        
    except Exception as e:
        return f"Lỗi khi đọc PPTX: {str(e)}"


def extract_excel_text(file):
    """Extract text from Excel file"""
    try:
        file.seek(0)
        
        # Try with pandas first
        try:
            # Read all sheets
            xl_file = pd.ExcelFile(file)
            text = ""
            
            for sheet_name in xl_file.sheet_names:
                text += f"\n--- Sheet: {sheet_name} ---\n"
                df = pd.read_excel(xl_file, sheet_name=sheet_name)
                
                # Convert to string, handling NaN values
                df_text = df.fillna('').astype(str)
                text += df_text.to_string(index=False) + "\n"
            
            return text.strip() if text.strip() else "File Excel này không chứa dữ liệu có thể đọc được."
            
        except Exception as e:
            # Fallback to openpyxl
            file.seek(0)
            wb = openpyxl.load_workbook(file)
            text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                        else:
                            row_text.append("")
                    if any(cell.strip() for cell in row_text):  # Skip empty rows
                        text += " | ".join(row_text) + "\n"
            
            return text.strip() if text.strip() else "File Excel này không chứa dữ liệu có thể đọc được."
            
    except Exception as e:
        return f"Lỗi khi đọc Excel: {str(e)}"


def extract_text_file(file):
    """Extract text from text files"""
    try:
        file.seek(0)
        
        # Try different encodings
        encodings = ['utf-8', 'utf-16', 'latin1', 'cp1252']
        
        for encoding in encodings:
            try:
                file.seek(0)
                content = file.read()
                if isinstance(content, bytes):
                    text = content.decode(encoding)
                else:
                    text = content
                
                return text.strip() if text.strip() else "File text này trống."
                
            except UnicodeDecodeError:
                continue
        
        return "Không thể đọc file text với các encoding thông dụng."
        
    except Exception as e:
        return f"Lỗi khi đọc text file: {str(e)}"


def image_to_base64(image_file):
    """Convert image file to base64 for Gemini API"""
    try:
        if hasattr(image_file, 'read'):
            image_data = image_file.read()
        else:
            with open(image_file, 'rb') as f:
                image_data = f.read()
        
        # Get image format
        image = Image.open(io.BytesIO(image_data))
        format_lower = image.format.lower() if image.format else 'jpeg'
        
        base64_string = base64.b64encode(image_data).decode('utf-8')
        
        return {
            'mime_type': f'image/{format_lower}',
            'data': base64_string
        }
    except Exception as e:
        print(f"Error converting image to base64: {e}")
        return None


def call_gemini_api(messages, image_data=None):
    """Call Gemini API with conversation history and optional image"""
    try:
        headers = {
            'Content-Type': 'application/json',
        }
        
        # Prepare contents for API
        contents = []
        
        # Add conversation history
        for msg in messages:
            if msg['role'] == 'system':
                continue  # Gemini doesn't use system role, we'll include it in prompt
                
            content = {
                "role": "user" if msg['role'] == 'user' else "model",
                "parts": [{"text": msg['content']}]
            }
            contents.append(content)
        
        # Add image if provided
        if image_data and contents:
            # Add image to the last user message
            for content in reversed(contents):
                if content['role'] == 'user':
                    content['parts'].append({
                        "inline_data": image_data
                    })
                    break
        
        payload = {
            "contents": contents,
            "generationConfig": {
                "temperature": 0.7,
                "topK": 40,
                "topP": 0.95,
                "maxOutputTokens": 8192,
            },
            "safetySettings": [
                {
                    "category": "HARM_CATEGORY_HARASSMENT",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                },
                {
                    "category": "HARM_CATEGORY_HATE_SPEECH",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                },
                {
                    "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                },
                {
                    "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
                    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
                }
            ]
        }
        
        response = requests.post(
            f"{GEMINI_URL}?key={GEMINI_API_KEY}",
            headers=headers,
            data=json.dumps(payload),
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            if 'candidates' in result and len(result['candidates']) > 0:
                content = result['candidates'][0]['content']['parts'][0]['text']
                return {
                    'success': True,
                    'content': content,
                    'usage': result.get('usageMetadata', {})
                }
            else:
                return {
                    'success': False,
                    'error': 'No response from AI'
                }
        else:
            return {
                'success': False,
                'error': f'API Error: {response.status_code} - {response.text}'
            }
            
    except Exception as e:
        return {
            'success': False,
            'error': f'Exception: {str(e)}'
        }


# 1. SỬA VIEW FUNCTION
@login_required
def ai_image_solver_view(request):
    """Main page for AI Image Solver"""
    # Debug: In ra console để check
    if request.user.is_authenticated:
        all_solutions = AIImageSolution.objects.filter(user=request.user).order_by('-created_at')
    else:
        all_solutions = AIImageSolution.objects.none()
    print(f"Total solutions for user {request.user.id}: {all_solutions.count()}")
    
    for solution in all_solutions[:5]:  # Debug 5 solutions đầu
        print(f"Solution {solution.id}: {solution.title}, type: {solution.solution_type}, image_url: {solution.image_url}")
    
    # Get user's recent solutions for history - bao gồm tất cả loại
    recent_solutions = all_solutions[:20]
    
    # Get active conversations
    active_conversations = AIConversation.objects.filter(
        user=request.user,
        is_active=True
    ).order_by('-updated_at')[:5]
    
    context = {
        'recent_solutions': recent_solutions,
        'active_conversations': active_conversations,
    }
    
    return render(request, 'ai/image_solver.html', context)

@login_required
@csrf_exempt
@require_http_methods(["POST"])
def ai_solve_image_api(request):
    """API endpoint to process image and get AI solution"""
    try:
        start_time = time.time()
        
        # Debug logs
        print(f"Request FILES: {request.FILES}")
        print(f"Request POST: {request.POST}")
        
        # Get uploaded image
        image_file = request.FILES.get('image')
        if not image_file:
            print("No image file found in request")
            return JsonResponse({
                'success': False,
                'error': 'No image provided'
            })
        
        # Get optional conversation ID and user question
        conversation_id = request.POST.get('conversation_id')
        user_question = request.POST.get('question', '')
        
        print(f"Conversation ID: {conversation_id}")
        print(f"User question: {user_question}")
        
        # FIXED: Tạo hoặc lấy conversation NGAY TỪ ĐẦU
        conversation = None
        if conversation_id:
            print("Looking for existing conversation:", conversation_id)
            try:
                conversation = AIConversation.objects.get(
                    id=conversation_id,
                    user=request.user
                )
                print("Found existing conversation:", conversation.id, conversation.title)
            except AIConversation.DoesNotExist:
                print("Conversation not found, will create new one")
                conversation = None
        # THÊM DEBUG LOGS VÀO ĐÂY:
        if conversation:
            print(f"=== DEBUG CONVERSATION {conversation.id} ===")
            print(f"Conversation title: {conversation.title}")
            print(f"Created at: {conversation.created_at}")
            print(f"Has solution: {bool(conversation.image_solution)}")
            
            # Check tất cả messages trong conversation này
            all_msgs = AIConversationMessage.objects.filter(conversation=conversation)
            print(f"Total messages in conversation: {all_msgs.count()}")
            
            for i, msg in enumerate(all_msgs.order_by('created_at')):
                print(f"  Message {i+1}: {msg.role} - {msg.content[:100]}...")
                print(f"              Created: {msg.created_at}")
            print("=== END DEBUG ===")
                # Nếu chưa có conversation, tạo mới NGAY
        if not conversation:
            print("Creating new IMAGE conversation early...")
            conversation = AIConversation.objects.create(
                user=request.user,
                title=f"Image Analysis - {time.strftime('%d/%m/%Y %H:%M')}",
                image_solution=None  # Sẽ update sau
            )
            print(f"Created conversation ID: {conversation.id}")
        
        # Validate image file
        if image_file.size == 0:
            print("Image file is empty")
            return JsonResponse({
                'success': False,
                'error': 'Uploaded file is empty'
            })
        
        if image_file.size > 10 * 1024 * 1024:  # 10MB
            return JsonResponse({
                'success': False,
                'error': 'File too large. Maximum size is 10MB.'
            })
        
        # Read and validate file content
        try:
            image_file.seek(0)
            file_content = image_file.read()
            
            if not file_content:
                return JsonResponse({
                    'success': False,
                    'error': 'File content is empty'
                })
            
            print(f"File content length: {len(file_content)}")
            image_file.seek(0)
            
        except Exception as e:
            print(f"Error reading file: {e}")
            return JsonResponse({
                'success': False,
                'error': f'Error reading file: {str(e)}'
            })
        
        # Convert image to base64
        try:
            image_data = image_to_base64(image_file)
            if not image_data:
                return JsonResponse({
                    'success': False,
                    'error': 'Failed to process image - conversion error'
                })
            print("Image converted to base64 successfully")
            
        except Exception as e:
            print(f"Error converting image: {e}")
            return JsonResponse({
                'success': False,
                'error': f'Failed to process image: {str(e)}'
            })
        
        # Prepare conversation messages (lấy từ conversation đã tạo)
        messages = [
            {
                'role': 'system',
                'content': '''Bạn là một AI assistant chuyên giải bài tập và trả lời câu hỏi từ hình ảnh. 
                Nhiệm vụ của bạn:
                1. Đọc và hiểu nội dung trong hình ảnh (văn bản, công thức, biểu đồ...)
                2. Xác định các câu hỏi hoặc bài tập cần giải
                3. Cung cấp lời giải chi tiết, từng bước
                4. Giải thích rõ ràng bằng tiếng Việt
                5. Nếu có nhiều câu hỏi, giải từng câu một cách có thứ tự
                
                Định dạng trả lời:
                **📖 Nội dung đã đọc được:**
                [Tóm tắt nội dung trong ảnh]
                
                **❓ Câu hỏi/Bài tập:**
                [Liệt kê các câu hỏi]
                
                **✅ Lời giải:**
                [Giải chi tiết từng bước]
                '''
            }
        ]
        
        # Add conversation history if exists
        conversation_messages = AIConversationMessage.objects.filter(
            conversation=conversation
        ).order_by('-created_at')[:10]
        conversation_messages = list(reversed(conversation_messages))
        
        for msg in conversation_messages:
            messages.append({
                'role': msg.role,
                'content': msg.content
            })
        print(f"Added {len(conversation_messages)} conversation messages")
        
        # Add current user message
        current_message = f"Hãy phân tích hình ảnh này và giải các bài tập/câu hỏi trong đó."
        if user_question:
            current_message += f" Câu hỏi cụ thể: {user_question}"
            
        messages.append({
            'role': 'user',
            'content': current_message
        })
        
        # FIXED: Lưu user message NGAY VÀO DATABASE (trước khi call API)
        print("Saving user message to conversation...")
        user_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='user',
            content=current_message
        )
        print(f"Saved user message ID: {user_msg.id}")
        
        print(f"Calling Gemini API with {len(messages)} messages")
        
        # Call Gemini API
        api_response = call_gemini_api(messages, image_data)
        
        if not api_response['success']:
            print(f"Gemini API error: {api_response['error']}")
            # FIXED: Vẫn có conversation và user message đã lưu
            return JsonResponse({
                'success': False,
                'error': api_response['error'],
                'conversation': {
                    'id': conversation.id,
                    'title': conversation.title,
                }
            })
        
        processing_time = int((time.time() - start_time) * 1000)
        ai_content = api_response['content']
        
        print(f"AI response received, processing time: {processing_time}ms")
        
        # Upload image to Cloudinary
        try:
            image_file.seek(0)
            from cloudinary import uploader
            upload_result = uploader.upload(
                image_file,
                folder="ai_images/",
                public_id=f"ai_solution_{int(time.time())}_{request.user.id}",
                resource_type="image"
            )
            print(f"Cloudinary upload successful: {upload_result['public_id']}")
            
        except Exception as e:
            print(f"Cloudinary upload error: {e}")
            # FIXED: Vẫn lưu AI response vào conversation
            ai_msg = AIConversationMessage.objects.create(
                conversation=conversation,
                role='assistant',
                content=f"Lỗi upload hình ảnh: {str(e)}\n\nNhưng đây là phân tích của AI:\n\n{ai_content}",
                tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
                response_time=processing_time
            )
            
            return JsonResponse({
                'success': False,
                'error': f'Failed to upload image: {str(e)}',
                'conversation': {
                    'id': conversation.id,
                    'title': conversation.title,
                }
            })
        
        # Create AIImageSolution
        try:
            solution = AIImageSolution.objects.create(
                user=request.user,
                image_url=upload_result['public_id'],
                original_filename=image_file.name,
                ai_solution=ai_content,
                processing_time=processing_time,
                title=f"AI Image Solution - {time.strftime('%Y-%m-%d %H:%M')}",
                solution_type='image'  # Set explicit type
            )
            print(f"AIImageSolution created: {solution.id}")
            
            # FIXED: Update conversation với solution
            conversation.image_solution = solution
            conversation.title = f"Image Analysis - {solution.title}"
            conversation.save()
            print(f"Updated conversation {conversation.id} with solution {solution.id}")
            
        except Exception as e:
            print(f"Error creating solution: {e}")
            # FIXED: Vẫn lưu AI response vào conversation
            ai_msg = AIConversationMessage.objects.create(
                conversation=conversation,
                role='assistant',
                content=ai_content,
                tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
                response_time=processing_time
            )
            
            return JsonResponse({
                'success': False,
                'error': f'Failed to save solution: {str(e)}',
                'conversation': {
                    'id': conversation.id,
                    'title': conversation.title,
                }
            })
        
        # FIXED: Lưu AI response message
        print("Saving AI response message...")
        ai_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=ai_content,
            tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
            response_time=processing_time
        )
        print(f"Saved AI message ID: {ai_msg.id}")
        
        # Update conversation timestamp
        conversation.save()
        
        return JsonResponse({
            'success': True,
            'solution': {
                'id': solution.id,
                'title': solution.title,
                'ai_solution': ai_content,
                'image_url': upload_result['secure_url'],
                'processing_time': processing_time,
                'created_at': solution.created_at.isoformat(),
            },
            'conversation': {
                'id': conversation.id,
                'title': conversation.title,
            }
        })
        
    except Exception as e:
        print(f"Unexpected error in ai_solve_image_api: {e}")
        import traceback
        print(traceback.format_exc())
        
        return JsonResponse({
            'success': False,
            'error': f'Server error: {str(e)}'
        })


@login_required
@csrf_exempt
@require_http_methods(["POST"])
def ai_solve_file_api(request):
    """API endpoint to process document files and get AI solution"""
    try:
        start_time = time.time()
        
        print(f"Request FILES: {request.FILES}")
        print(f"Request POST: {request.POST}")
        
        # Get uploaded file
        file_upload = request.FILES.get('file')
        if not file_upload:
            return JsonResponse({
                'success': False,
                'error': 'No file provided'
            })
        
        # Get optional parameters
        conversation_id = request.POST.get('conversation_id')
        user_question = request.POST.get('question', '')
        
        print(f"File name: {file_upload.name}")
        print(f"Conversation ID: {conversation_id}")
        print(f"User question: {user_question}")
        
        # FIXED: Tạo hoặc lấy conversation NGAY TỪ ĐẦU
        conversation = None
        if conversation_id:
            try:
                conversation = AIConversation.objects.get(
                    id=conversation_id,
                    user=request.user
                )
                print(f"Found existing conversation: {conversation.id}")
            except AIConversation.DoesNotExist:
                print("Conversation not found, will create new one")
                conversation = None
        
        # Nếu chưa có conversation, tạo mới NGAY
        if not conversation:
            print("Creating new FILE conversation early...")
            conversation = AIConversation.objects.create(
                user=request.user,
                title=f"Document Analysis - {file_upload.name}",
                image_solution=None  # Sẽ update sau
            )
            print(f"Created conversation ID: {conversation.id}")
        
        # Validate file
        if file_upload.size > 20 * 1024 * 1024:
            return JsonResponse({
                'success': False,
                'error': 'File too large. Maximum size is 20MB.'
            })
        
        file_extension = os.path.splitext(file_upload.name)[1].lower()
        if file_extension not in ['.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.txt', '.csv']:
            return JsonResponse({
                'success': False,
                'error': f'File type {file_extension} not supported. Supported types: PDF, DOCX, PPTX, XLSX, XLS, TXT, CSV'
            })
        
        # Extract text from file
        try:
            extracted_text = extract_text_from_file(file_upload)
            print(f"Extracted text length: {len(extracted_text)} characters")
            
            if not extracted_text or len(extracted_text.strip()) < 10:
                return JsonResponse({
                    'success': False,
                    'error': 'Could not extract meaningful content from the file'
                })
                
        except Exception as e:
            print(f"Error extracting text: {e}")
            return JsonResponse({
                'success': False,
                'error': f'Failed to extract text from file: {str(e)}'
            })
        
        # Prepare conversation messages
        messages = [
            {
                'role': 'system',
                'content': '''Bạn là một AI assistant chuyên phân tích và giải đáp từ các tài liệu văn bản.
                Nhiệm vụ của bạn:
                1. Đọc và hiểu nội dung trong tài liệu
                2. Trả lời câu hỏi dựa trên nội dung tài liệu
                3. Giải thích chi tiết và rõ ràng bằng tiếng Việt
                4. Nếu có bài tập, hãy giải từng bước
                5. Tóm tắt nội dung chính nếu được yêu cầu
                
                Định dạng trả lời:
                **📄 Tóm tắt nội dung tài liệu:**
                [Tóm tắt ngắn gọn nội dung chính]
                
                **❓ Câu hỏi/Yêu cầu:**
                [Câu hỏi của người dùng]
                
                **✅ Trả lời:**
                [Trả lời chi tiết dựa trên nội dung tài liệu]
                '''
            }
        ]
        
        # Add conversation history if exists
        conversation_messages = AIConversationMessage.objects.filter(
            conversation=conversation
        ).order_by('-created_at')[:10]
        conversation_messages = list(reversed(conversation_messages))
        
        for msg in conversation_messages:
            messages.append({
                'role': msg.role,
                'content': msg.content
            })
        
        # Truncate extracted text if too long
        if len(extracted_text) > 15000:
            extracted_text = extracted_text[:15000] + "\n\n[Nội dung bị cắt do quá dài...]"
        
        # Add current user message with file content
        current_message = f"Đây là nội dung từ file '{file_upload.name}':\n\n{extracted_text}\n\n"
        
        if user_question:
            current_message += f"Câu hỏi của tôi: {user_question}"
        else:
            current_message += "Hãy tóm tắt nội dung chính và trả lời bất kỳ câu hỏi nào có trong tài liệu."
            
        messages.append({
            'role': 'user',
            'content': current_message
        })
        
        # FIXED: Lưu user message NGAY VÀO DATABASE (trước khi call API)
        print("Saving user message to conversation...")
        user_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='user',
            content=f"Uploaded file: {file_upload.name}\n\nQuestion: {user_question or 'Analyze this document'}"
        )
        print(f"Saved user message ID: {user_msg.id}")
        
        print(f"Calling Gemini API with {len(messages)} messages")
        
        # Call Gemini API (no image data)
        api_response = call_gemini_api(messages, image_data=None)
        
        if not api_response['success']:
            print(f"Gemini API error: {api_response['error']}")
            # FIXED: Vẫn có conversation và user message đã lưu
            return JsonResponse({
                'success': False,
                'error': api_response['error'],
                'conversation': {
                    'id': conversation.id,
                    'title': conversation.title,
                }
            })
        
        processing_time = int((time.time() - start_time) * 1000)
        ai_content = api_response['content']
        
        # Upload file to Cloudinary
        try:
            file_upload.seek(0)
            from cloudinary import uploader
            upload_result = uploader.upload(
                file_upload,
                folder="ai_documents/",
                public_id=f"ai_doc_{int(time.time())}_{request.user.id}",
                resource_type="raw"
            )
            print(f"File uploaded to Cloudinary: {upload_result['public_id']}")
            
        except Exception as e:
            print(f"Cloudinary upload error: {e}")
            # FIXED: Vẫn lưu AI response vào conversation
            ai_msg = AIConversationMessage.objects.create(
                conversation=conversation,
                role='assistant',
                content=f"Lỗi upload file: {str(e)}\n\nNhưng đây là phân tích của AI:\n\n{ai_content}",
                tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
                response_time=processing_time
            )
            
            return JsonResponse({
                'success': False,
                'error': f'Failed to upload file: {str(e)}',
                'conversation': {
                    'id': conversation.id,
                    'title': conversation.title,
                }
            })
        
        # Create AIImageSolution (reusing model for document solutions)
        try:
            solution = AIImageSolution.objects.create(
                user=request.user,
                image_url=upload_result['public_id'],
                original_filename=file_upload.name,
                extracted_text=extracted_text[:5000] if len(extracted_text) > 5000 else extracted_text,
                ai_solution=ai_content,
                processing_time=processing_time,
                title=f"Document Analysis - {file_upload.name}",
                solution_type='document',
                document_type=file_extension[1:],
                file_size=file_upload.size
            )
            print(f"AIImageSolution created: {solution.id}")
            
            # FIXED: Update conversation với solution
            conversation.image_solution = solution
            conversation.title = f"Document Analysis - {file_upload.name}"
            conversation.save()
            print(f"Updated conversation {conversation.id} with solution {solution.id}")
            
        except Exception as e:
            print(f"Error creating solution: {e}")
            # FIXED: Vẫn lưu AI response vào conversation
            ai_msg = AIConversationMessage.objects.create(
                conversation=conversation,
                role='assistant',
                content=ai_content,
                tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
                response_time=processing_time
            )
            
            return JsonResponse({
                'success': False,
                'error': f'Failed to save solution: {str(e)}',
                'conversation': {
                    'id': conversation.id,
                    'title': conversation.title,
                }
            })
        
        # FIXED: Lưu AI response message
        print("Saving AI response message...")
        ai_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=ai_content,
            tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
            response_time=processing_time
        )
        print(f"Saved AI message ID: {ai_msg.id}")
        
        # Update conversation timestamp
        conversation.save()
        
        return JsonResponse({
            'success': True,
            'solution': {
                'id': solution.id,
                'title': solution.title,
                'ai_solution': ai_content,
                'file_url': upload_result['secure_url'],
                'file_name': file_upload.name,
                'processing_time': processing_time,
                'created_at': solution.created_at.isoformat(),
            },
            'conversation': {
                'id': conversation.id,
                'title': conversation.title,
            }
        })
        
    except Exception as e:
        print(f"Unexpected error in ai_solve_file_api: {e}")
        import traceback
        print(traceback.format_exc())
        
        return JsonResponse({
            'success': False,
            'error': f'Server error: {str(e)}'
        })



@login_required
@csrf_exempt
@require_http_methods(["POST"])
def ai_search_documents_api(request):
    """API để tìm kiếm tài liệu theo yêu cầu của user"""
    try:
        query = request.POST.get('query', '').strip()
        if not query:
            return JsonResponse({
                'success': False,
                'error': 'Missing search query'
            })
        
        documents = search_documents_for_ai(query, request.user, limit=10)
        
        return JsonResponse({
            'success': True,
            'documents': documents,
            'count': len(documents)
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': f'Search error: {str(e)}'
        })

# NEW: API endpoint for manual chat room search  
@login_required
@csrf_exempt
@require_http_methods(["POST"])
def ai_search_chat_rooms_api(request):
    """API để tìm kiếm phòng chat theo yêu cầu của user"""
    try:
        query = request.POST.get('query', '').strip()
        if not query:
            return JsonResponse({
                'success': False,
                'error': 'Missing search query'
            })
        
        chat_rooms = search_chat_rooms_for_ai(query, request.user, limit=10)
        
        return JsonResponse({
            'success': True,
            'chat_rooms': chat_rooms,
            'count': len(chat_rooms)
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': f'Search error: {str(e)}'
        })

@login_required
@csrf_exempt
@require_http_methods(["POST"])
def ai_search_courses_api(request):
    """API để tìm kiếm khóa học lập trình theo yêu cầu của user"""
    try:
        query = request.POST.get('query', '').strip()
        if not query:
            return JsonResponse({
                'success': False,
                'error': 'Missing search query'
            })
        
        # Import CodeCourse model
        from .models import CodeCourse
        
        print(f"=== DEBUG SEARCH COURSES ===")
        print(f"Query: '{query}'")
        print(f"User: {request.user.username}")
        
        # Search courses with fuzzy matching
        courses = CodeCourse.objects.filter(
            Q(status='published') &
            (Q(title__icontains=query) |
             Q(description__icontains=query) |
             Q(language__display_name__icontains=query) |
             Q(language__name__icontains=query) |
             Q(difficulty__icontains=query))
        ).select_related('language', 'created_by', 'university').order_by('-enrollment_count', '-rating_average')[:10]
        
        print(f"Found courses: {courses.count()}")
        
        # Fallback if no results
        if courses.count() == 0:
            print("No results with query, trying fallback...")
            courses = CodeCourse.objects.filter(
                status='published'
            ).select_related('language', 'created_by', 'university').order_by('-enrollment_count', '-rating_average')[:10]
            print(f"Fallback results: {courses.count()}")
        
        # Format results
        results = []
        for course in courses:
            try:
                result_item = {
                    'id': course.id,
                    'title': course.title,
                    'description': course.description[:200] if course.description else '',
                    'language': course.language.display_name,
                    'difficulty': course.get_difficulty_display() if hasattr(course, 'get_difficulty_display') else course.difficulty,
                    'enrollment_count': course.enrollment_count,
                    'rating': float(course.rating_average) if course.rating_average else 0,
                    'url': f"/code/courses/{course.slug}/"
                }
                results.append(result_item)
                print(f"Formatted course: {result_item['title']}")
            except Exception as e:
                print(f"Error formatting course {course.id}: {e}")
                continue
        
        print(f"Formatted results: {len(results)}")
        print("=== END DEBUG ===")
        
        return JsonResponse({
            'success': True,
            'courses': results,
            'count': len(results)
        })
        
    except Exception as e:
        print(f"Course search error: {e}")
        import traceback
        traceback.print_exc()
        return JsonResponse({
            'success': False,
            'error': f'Search error: {str(e)}'
        })

@login_required
@csrf_exempt
@require_http_methods(["POST"])
def ai_export_solutions(request):
    """Export selected AI solutions to PDF"""
    try:
        import csv
        from io import StringIO
        
        data = json.loads(request.body)
        solution_ids = data.get('solution_ids', [])
        
        if not solution_ids:
            return JsonResponse({
                'success': False,
                'error': 'No solutions selected'
            })
        
        # Get solutions
        solutions = AIImageSolution.objects.filter(
            id__in=solution_ids,
            user=request.user
        ).order_by('-created_at')
        
        if not solutions.exists():
            return JsonResponse({
                'success': False,
                'error': 'No solutions found'
            })
        
        # Create CSV file
        output = StringIO()
        writer = csv.writer(output)
        
        # Write headers
        writer.writerow(['ID', 'Tiêu đề', 'Loại', 'Ngày tạo', 'Thời gian xử lý (ms)', 'Lượt xem', 'Lượt thích', 'Nội dung AI'])
        
        # Write data
        for solution in solutions:
            solution_type = 'Hình ảnh' if solution.solution_type == 'image' else ('Tài liệu' if solution.solution_type == 'document' else 'Text Chat')
            writer.writerow([
                solution.id,
                solution.title,
                solution_type,
                solution.created_at.strftime('%d/%m/%Y %H:%M'),
                solution.processing_time or 0,
                solution.view_count,
                solution.like_count,
                solution.ai_solution[:500] + '...' if len(solution.ai_solution) > 500 else solution.ai_solution
            ])
        
        # Create response
        response = HttpResponse(output.getvalue(), content_type='text/csv')
        response['Content-Disposition'] = f'attachment; filename="ai_solutions_{timezone.now().strftime("%Y%m%d_%H%M%S")}.csv"'
        
        return response
        
    except Exception as e:
        print(f"Export error: {e}")
        import traceback
        traceback.print_exc()
        return JsonResponse({
            'success': False,
            'error': f'Export error: {str(e)}'
        })

# Keep all other existing functions unchanged...
# (Include all the existing functions like ai_solve_image_api, ai_solve_file_api, etc.)
@login_required
@csrf_exempt  
@require_http_methods(["POST"])
def ai_continue_conversation_api(request):
    """API để tiếp tục cuộc trò chuyện với AI (enhanced)"""
    try:
        start_time = time.time()
        
        conversation_id = request.POST.get('conversation_id')
        user_message = request.POST.get('message', '').strip()
        
        if not conversation_id or not user_message:
            return JsonResponse({
                'success': False,
                'error': 'Missing conversation ID or message'
            })
        
        # Get conversation
        try:
            conversation = AIConversation.objects.get(
                id=conversation_id,
                user=request.user
            )
        except AIConversation.DoesNotExist:
            return JsonResponse({
                'success': False,
                'error': 'Conversation not found'
            })
        
        # Get conversation history
        conversation_messages = AIConversationMessage.objects.filter(
            conversation=conversation
        ).order_by('-created_at')[:15]
        conversation_messages = list(reversed(conversation_messages))
        
        # Prepare messages for API
        messages = []
        
        # Add conversation history
        for msg in conversation_messages:
            messages.append({
                'role': msg.role,
                'content': msg.content
            })
        
        # Add current user message
        messages.append({
            'role': 'user',
            'content': user_message
        })
        
        # ENHANCED: Call API với database context
        api_response = call_gemini_api_enhanced(messages, user=request.user)
        
        if not api_response['success']:
            return JsonResponse({
                'success': False,
                'error': api_response['error']
            })
        
        processing_time = int((time.time() - start_time) * 1000)
        ai_response = api_response['content']
        
        # Save messages
        user_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='user',
            content=user_message
        )
        
        ai_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=ai_response,
            tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
            response_time=processing_time
        )
        
        # Update conversation timestamp
        conversation.save()
        
        return JsonResponse({
            'success': True,
            'response': {
                'content': ai_response,
                'processing_time': processing_time,
                'created_at': ai_msg.created_at.isoformat(),
            }
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': f'Server error: {str(e)}'
        })

@login_required
@csrf_exempt
@require_http_methods(["POST"])
def ai_text_chat_api(request):
    """API để chat với AI bằng text thuần túy với database integration"""
    print("=== AI TEXT CHAT API CALLED (ENHANCED) ===")
    
    try:
    
        start_time = time.time()
    
        user_message = request.POST.get('message', '').strip()
        conversation_id = request.POST.get('conversation_id')
        
        print("User message:", repr(user_message))
        print("User:", request.user.username)
        
        # Nếu không có message, tạo conversation mới và trả về
        if not user_message:
            if not conversation_id:
                conversation = AIConversation.objects.create(
                    user=request.user,
                    title=f"Chat AI - {time.strftime('%d/%m/%Y %H:%M')}",
                    image_solution=None
                )
                return JsonResponse({
                    'success': True,
                    'conversation_id': conversation.id,
                    'message': 'Conversation created'
                })
            else:
                return JsonResponse({
                    'success': True,
                    'conversation_id': conversation_id,
                    'message': 'Ready to chat'
                })
        
        # Tạo hoặc lấy conversation
        conversation = None
        is_new_conversation = False
        
        if conversation_id:
            try:
                conversation = AIConversation.objects.get(
                    id=conversation_id,
                    user=request.user
                )
            except AIConversation.DoesNotExist:
                conversation = None
        
        if not conversation:
            is_new_conversation = True
            conversation = AIConversation.objects.create(
                user=request.user,
                title=f"Chat AI - {time.strftime('%d/%m/%Y %H:%M')}",
                image_solution=None
            )
        
        # Lấy lịch sử conversation
        conversation_messages = AIConversationMessage.objects.filter(
            conversation=conversation
        ).order_by('created_at')
        
        messages = []
        
        # Thêm lịch sử conversation
        recent_messages = conversation_messages[-10:] if len(conversation_messages) > 10 else conversation_messages
        
        for msg in recent_messages:
            messages.append({
                'role': msg.role,
                'content': msg.content
            })
        
        # Thêm tin nhắn hiện tại
        messages.append({
            'role': 'user',
            'content': user_message
        })
        
        # ENHANCED: Gọi API với database integration
        # TODO: Add rate limiting here to prevent API overload
        print("Calling enhanced Gemini API with DB context...")
        api_response = call_gemini_api_enhanced(messages, image_data=None, user=request.user)
        
        if not api_response['success']:
            return JsonResponse({
                'success': False,
                'error': api_response['error']
            })
        
        processing_time = int((time.time() - start_time) * 1000)
        ai_response = api_response['content']
        
        # Tạo solution cho conversation chưa có solution
        if not conversation.image_solution:
            try:
                solution = AIImageSolution.objects.create(
                    user=request.user,
                    title=f"Text Chat - {time.strftime('%d/%m/%Y %H:%M')}",
                    ai_solution=ai_response,
                    solution_type='text_chat',
                    processing_time=processing_time,
                    original_filename='text_chat.txt'
                )
                
                conversation.image_solution = solution
                conversation.save()
                
            except Exception as e:
                print("ERROR creating enhanced text chat solution:", str(e))
        
        # Lưu messages vào database
        user_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='user',
            content=user_message
        )
        
        ai_msg = AIConversationMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=ai_response,
            tokens_used=api_response.get('usage', {}).get('totalTokenCount', 0),
            response_time=processing_time
        )
        
        conversation.save()
        
        return JsonResponse({
            'success': True,
            'conversation_id': conversation.id,
            'response': {
                'content': ai_response,
                'processing_time': processing_time,
                'created_at': ai_msg.created_at.isoformat(),
            }
        })
        
    except Exception as e:
        print("=== UNEXPECTED ERROR IN ENHANCED CHAT ===")
        print("Error:", str(e))
        import traceback
        traceback.print_exc()
        
        return JsonResponse({
            'success': False,
            'error': f'Server error: {str(e)}'
        })

@login_required
def ai_solution_detail_view(request, solution_id):
    """View chi tiết một AI solution với conversation history"""
    solution = get_object_or_404(AIImageSolution, id=solution_id, user=request.user)
    
    # Increment view count
    solution.view_count += 1
    solution.save()
    
    # Get conversations related to this solution
    conversations = AIConversation.objects.filter(image_solution=solution)
    
    # Get all messages for these conversations - Fixed structure
    conversations_with_messages = []
    for conv in conversations:
        messages = AIConversationMessage.objects.filter(
            conversation=conv
        ).order_by('created_at')
        conversations_with_messages.append({
            'conversation': conv,
            'messages': messages
        })
    
    context = {
        'solution': solution,
        'conversations': conversations,
        'conversations_with_messages': conversations_with_messages,  # Fixed key name
    }
    
    return render(request, 'ai/solution_detail.html', context)


def ai_solutions_history_view(request):
    """View lịch sử các AI solutions"""
    if request.user.is_authenticated:
        solutions = AIImageSolution.objects.filter(user=request.user).order_by('-created_at')
    else:
        solutions = AIImageSolution.objects.none()
    
    # Pagination if needed
    from django.core.paginator import Paginator
    paginator = Paginator(solutions, 20)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)
    
    context = {
        'solutions': page_obj,
    }
    
    return render(request, 'ai/solutions_history.html', context)


@login_required
@csrf_exempt
@require_http_methods(["POST"])
def user_report_api(request):
    """API để người dùng báo cáo bài giải"""
    try:
        solution_id = request.POST.get('solution_id')
        reason = request.POST.get('reason')
        description = request.POST.get('description', '')
        
        if not solution_id or not reason:
            return JsonResponse({
                'success': False,
                'error': 'Thiếu thông tin bắt buộc'
            })
        
        # Kiểm tra solution có tồn tại không
        try:
            solution = AIImageSolution.objects.get(id=solution_id)
        except AIImageSolution.DoesNotExist:
            return JsonResponse({
                'success': False,
                'error': 'Bài giải không tồn tại'
            })
        
        # Gửi email báo cáo cho admin
        from django.core.mail import send_mail
        from django.conf import settings
        
        subject = f'Báo cáo bài giải AI - ID: {solution_id}'
        message = f"""
        Người báo cáo: {request.user.username} ({request.user.email})
        Bài giải ID: {solution_id}
        Tiêu đề: {solution.title}
        Lý do: {reason}
        Mô tả: {description}
        
        Link bài giải: {request.build_absolute_uri(f'/ai/solution/{solution_id}/')}
        """
        
        try:
            # Chỉ gửi email nếu đã cấu hình email settings
            if hasattr(settings, 'EMAIL_HOST') and settings.EMAIL_HOST:
                send_mail(
                    subject,
                    message,
                    settings.DEFAULT_FROM_EMAIL,
                    [settings.DEFAULT_FROM_EMAIL],  # Gửi cho admin
                    fail_silently=False,
                )
        except Exception as e:
            print(f"Email sending failed: {e}")
        
        # Log báo cáo
        print(f"Report submitted by {request.user.username} for solution {solution_id}: {reason}")
        
        return JsonResponse({
            'success': True,
            'message': 'Báo cáo đã được gửi thành công'
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': f'Lỗi server: {str(e)}'
        })
    

@login_required
@csrf_exempt  
@require_http_methods(["GET"])
def get_conversation_api(request, conversation_id):
    """API endpoint để lấy thông tin conversation và messages"""
    try:
        # Get conversation
        conversation = get_object_or_404(
            AIConversation, 
            id=conversation_id, 
            user=request.user
        )
        
        # Get messages
        messages = AIConversationMessage.objects.filter(
            conversation=conversation
        ).order_by('created_at')
        
        # Get solution info
        solution = conversation.image_solution
        
        # Prepare response data
        conversation_data = {
            'id': conversation.id,
            'title': conversation.title,
            'created_at': conversation.created_at.isoformat(),
            'updated_at': conversation.updated_at.isoformat(),
        }
        
        # Add solution info if exists
        if solution:
            conversation_data.update({
                'solution_type': solution.solution_type,
                'original_filename': solution.original_filename,
                'document_type': solution.document_type,
            })
            
            # Add image URL if exists and not text chat
            if solution.image_url and solution.solution_type != 'text_chat':
                from cloudinary import CloudinaryImage
                try:
                    # Get Cloudinary URL
                    if hasattr(solution.image_url, 'url'):
                        conversation_data['image_url'] = solution.image_url.url
                    else:
                        # Generate Cloudinary URL
                        img = CloudinaryImage(solution.image_url)
                        conversation_data['image_url'] = img.build_url()
                except:
                    conversation_data['image_url'] = None
        
        # Prepare messages data
        messages_data = []
        for msg in messages:
            msg_data = {
                'id': msg.id,
                'role': msg.role,
                'content': msg.content,
                'created_at': msg.created_at.isoformat(),
                'tokens_used': msg.tokens_used,
                'response_time': msg.response_time,
            }
            
            # Add image URL if exists
            if msg.image_url:
                try:
                    if hasattr(msg.image_url, 'url'):
                        msg_data['image_url'] = msg.image_url.url
                    else:
                        from cloudinary import CloudinaryImage
                        img = CloudinaryImage(msg.image_url)
                        msg_data['image_url'] = img.build_url()
                except:
                    msg_data['image_url'] = None
                    
            messages_data.append(msg_data)
        
        return JsonResponse({
            'success': True,
            'conversation': conversation_data,
            'messages': messages_data
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': f'Error loading conversation: {str(e)}'
        })
    

def delete_search_history_item_api(request):
    """API to delete a single search history item"""
    if not request.user.is_authenticated:
        return JsonResponse({'success': False, 'message': 'Vui lòng đăng nhập'}, status=401)
    
    if request.method == 'POST':
        import json
        data = json.loads(request.body)
        query = data.get('query')
        
        if not query:
            return JsonResponse({'success': False, 'message': 'Thiếu query'}, status=400)
        
        # Delete all search history entries for this query and user
        SearchHistory.objects.filter(user=request.user, query=query).delete()
        return JsonResponse({'success': True, 'message': 'Đã xóa lịch sử'})
    
    return JsonResponse({'success': False, 'message': 'Phương thức không hợp lệ'}, status=405)


def clear_search_history_api(request):
    """API to clear user's search history"""
    if not request.user.is_authenticated:
        return JsonResponse({'success': False, 'message': 'Vui lòng đăng nhập'}, status=401)
    
    if request.method == 'POST':
        # Delete all search history for this user
        SearchHistory.objects.filter(user=request.user).delete()
        return JsonResponse({'success': True, 'message': 'Đã xóa lịch sử tìm kiếm'})
    
    return JsonResponse({'success': False, 'message': 'Phương thức không hợp lệ'}, status=405)


def normalize_vietnamese(text):
    """Normalize Vietnamese text for fuzzy search"""
    from unidecode import unidecode
    return unidecode(text.lower()).strip()


def fuzzy_match(query, text, threshold=0.6):
    """Check if query fuzzy matches text using Levenshtein distance"""
    from Levenshtein import ratio
    
    query_norm = normalize_vietnamese(query)
    text_norm = normalize_vietnamese(text)
    
    # Check if query is substring
    if query_norm in text_norm:
        return True
    
    # Check Levenshtein ratio (0.6 = 60% similarity)
    words_in_text = text_norm.split()
    for word in words_in_text:
        if ratio(query_norm, word) >= threshold:
            return True
    
    # Check if all query words exist in text
    query_words = query_norm.split()
    if len(query_words) > 1:
        matches = sum(1 for qw in query_words if any(ratio(qw, tw) >= threshold for tw in words_in_text))
        if matches >= len(query_words) * 0.7:  # 70% of words match
            return True
    
    return False


def search_suggestions_api(request):
    """API for real-time search autocomplete suggestions with fuzzy matching"""
    query = request.GET.get('q', '').strip()
    normalized_query = normalize_vietnamese(query)
    
    if len(query) < 1:
        # Return search history and popular searches when no query
        suggestions = {
            'history': [],
            'popular': [],
            'results': []
        }
        
        # Get user's search history if authenticated
        if request.user.is_authenticated:
            recent_searches = SearchHistory.objects.filter(
                user=request.user
            ).values('query').distinct().order_by('-created_at')[:5]
            suggestions['history'] = [s['query'] for s in recent_searches]
        
        # Get popular searches (top 5 most searched)
        from django.db.models import Count
        popular_searches = SearchHistory.objects.values('query').annotate(
            search_count=Count('id')
        ).order_by('-search_count')[:5]
        suggestions['popular'] = [s['query'] for s in popular_searches]
        
        return JsonResponse(suggestions)
    
    # Real-time suggestions when typing
    suggestions = {
        'documents': [],
        'courses': [],
        'universities': []
    }
    
    from django.db.models import Q
    
    # Search documents with fuzzy matching
    documents_qs = Document.objects.filter(
        Q(status='approved') & Q(is_public=True)
    ).select_related('course', 'university')
    
    # Exact match first
    exact_docs = list(documents_qs.filter(
        Q(title__icontains=query) | Q(description__icontains=query)
    )[:5])
    
    # Fuzzy match if not enough results
    if len(exact_docs) < 3 and normalized_query:
        seen_ids = {d.id for d in exact_docs}
        all_docs = documents_qs.exclude(id__in=seen_ids)[:100]
        
        for doc in all_docs:
            if fuzzy_match(query, doc.title, threshold=0.5):
                exact_docs.append(doc)
                if len(exact_docs) >= 5:
                    break
    
    documents = exact_docs[:5]
    
    for doc in documents:
        suggestions['documents'].append({
            'id': doc.id,
            'title': doc.title,
            'type': 'document',
            'course': f"{doc.course.code} - {doc.course.name}" if doc.course else '',
            'university': doc.university.short_name or doc.university.name if doc.university else '',
            'url': f'/documents/{doc.id}/view/',
            'icon': 'fa-file-alt'
        })
    
    # Search courses with fuzzy matching
    courses_qs = Course.objects.filter(is_active=True).select_related('university')
    
    exact_courses = list(courses_qs.filter(
        Q(name__icontains=query) | Q(code__icontains=query)
    )[:5])
    
    if len(exact_courses) < 3 and normalized_query:
        seen_ids = {c.id for c in exact_courses}
        all_courses = courses_qs.exclude(id__in=seen_ids)[:100]
        
        for course in all_courses:
            if fuzzy_match(query, course.name, threshold=0.5):
                exact_courses.append(course)
                if len(exact_courses) >= 5:
                    break
    
    courses = exact_courses[:5]
    
    for course in courses:
        suggestions['courses'].append({
            'id': course.id,
            'title': f"{course.code} - {course.name}",
            'type': 'course',
            'university': course.university.short_name or course.university.name,
            'url': f'/documents/?course={course.id}',
            'icon': 'fa-book'
        })
    
    # Search universities with fuzzy matching
    unis_qs = University.objects.filter(is_active=True)
    
    exact_unis = list(unis_qs.filter(
        Q(name__icontains=query) | Q(short_name__icontains=query)
    )[:5])
    
    if len(exact_unis) < 3 and normalized_query:
        seen_ids = {u.id for u in exact_unis}
        all_unis = unis_qs.exclude(id__in=seen_ids)[:100]
        
        for uni in all_unis:
            if fuzzy_match(query, uni.name, threshold=0.5) or \
               (uni.short_name and fuzzy_match(query, uni.short_name, threshold=0.5)):
                exact_unis.append(uni)
                if len(exact_unis) >= 5:
                    break
    
    universities = exact_unis[:5]
    
    for uni in universities:
        suggestions['universities'].append({
            'id': uni.id,
            'title': uni.name,
            'subtitle': uni.short_name or '',
            'type': 'university',
            'url': f'/documents/?university={uni.id}',
            'icon': 'fa-university'
        })
    
    return JsonResponse(suggestions)


def documents_list(request):
    # Filter logic
    documents = Document.objects.filter(status='approved', is_public=True)
    
    # Apply filters from GET parameters
    if request.GET.get('university'):
        documents = documents.filter(university_id=request.GET.get('university'))
    
    if request.GET.get('course'):
        documents = documents.filter(course_id=request.GET.get('course'))
        
    # Add pagination, context
    paginator = Paginator(documents, 12)
    page = request.GET.get('page')
    documents = paginator.get_page(page)
    
    return render(request, 'documents/list.html', {
        'documents': documents,
        'universities': University.objects.filter(is_active=True),
        'courses': Course.objects.filter(is_active=True),
    })


@csrf_exempt
def google_callback(request):
    """Handle Google OAuth callback"""
    print("----- Bắt đầu xử lý callback -----")
    code = request.GET.get('code')
    print(f"Code nhận được: {code}")

    # TẠO REDIRECT URI CHÍNH XÁC
    if settings.DEBUG:
        redirect_uri = 'http://localhost:8000/accounts/google/login/callback/'
    else:
        redirect_uri = 'https://doan4-django.vercel.app/accounts/google/login/callback/'
    
    print(f"Redirect URI được sử dụng: {redirect_uri}")

    # Trao đổi code lấy token
    token_url = 'https://oauth2.googleapis.com/token'
    data = {
        'code': code,
        'client_id': settings.GOOGLE_CLIENT_ID,
        'client_secret': settings.GOOGLE_CLIENT_SECRET,
        'redirect_uri': redirect_uri,  # SỬA ĐÂY
        'grant_type': 'authorization_code'
    }

    print("----- Gửi yêu cầu lấy token tới Google -----")
    print(f"Data gửi tới Google: {data}")  # THÊM LOG ĐỂ DEBUG
    
    response = requests.post(token_url, data=data)
    print(f"Phản hồi từ Google (token): {response.status_code}, {response.text}")

    if not response.ok:
        print(f"Lỗi khi trao đổi code: {response.status_code}, {response.text}")
        return JsonResponse({'success': False, 'error': 'Failed to get token'}, status=400)

    token_data = response.json()
    access_token = token_data.get('access_token')
    print(f"Access token nhận được: {access_token}")

    # Lấy thông tin người dùng
    user_info_url = 'https://www.googleapis.com/oauth2/v3/userinfo'
    headers = {'Authorization': f'Bearer {access_token}'}
    user_response = requests.get(user_info_url, headers=headers)
    print(f"Phản hồi từ Google (user info): {user_response.status_code}, {user_response.text}")

    if not user_response.ok:
        print(f"Lỗi khi lấy thông tin người dùng: {user_response.status_code}, {user_response.text}")
        return JsonResponse({'success': False, 'error': 'Failed to get user info'}, status=400)

    user_info = user_response.json()
    print(f"Thông tin người dùng từ Google: {user_info}")
    email = user_info.get('email')
    google_id = user_info.get('sub')

    if not email:
        return JsonResponse({'success': False, 'error': 'Email not found in user info'}, status=400)

    # Get or create user
    user, created = User.objects.get_or_create(
        email=email,
        defaults={
            'username': email.split('@')[0],
            'first_name': user_info.get('given_name', ''),
            'last_name': user_info.get('family_name', ''),
            'google_id': google_id,
        }
    )
    
    # Update google_id if not set
    if not user.google_id:
        user.google_id = google_id
        user.save()
        
    print(f"User: {user}, Created: {created}")

    # Login user
    login(request, user)
    
    messages.success(request, f'Chào mừng {user.get_full_name() or user.username}!')
    return redirect('dashboard')


